import streamlit as st
import pandas as pd
import subprocess
import json
import urllib.request
import urllib.parse
import os
import datetime
import io
import shutil
import sys
import zipfile
import pypdf
import docx
import openpyxl
import re

# --- CONFIG ---
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://ollama:11434")
LITELLM_URL = os.getenv("LITELLM_URL", "http://litellm:4000/v1")
QDRANT_URL = "http://qdrant:6333"
SEARXNG_URL = os.getenv("SEARXNG_URL", "http://searxng:8080")
WORK_DIR = "/work/qa_reports" # Reports
CONSUME_DIR = "/consume" # Paperless Input
IATF_BASE_URL = os.getenv("IATF_BASE_URL", "http://localhost:3000")
PROCESS_MONITORING_URL = f"{IATF_BASE_URL}/products/process_monitoring_measurement"

# Subfolders
INGEST_DIR = os.path.join(CONSUME_DIR, "PFMEA_5WHY_FTA_etc")
WIP_DIR = os.path.join(CONSUME_DIR, "WIP")
KINDLE_DIR = os.path.join(CONSUME_DIR, "Kindle")

GEN_MODEL = os.getenv("OLLAMA_GEN_MODEL", "qwen3:14b")
EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")
FMEA_DEEP_MODEL = os.getenv("FMEA_DEEP_MODEL", os.getenv("OPENAI_MODEL", "openai/gemini-2.5-flash"))
FMEA_DEEP_API_KEY = os.getenv("FMEA_DEEP_API_KEY", os.getenv("OPENAI_API_KEY", ""))
WHYWHY_MODEL_GEMINI = os.getenv("WHYWHY_MODEL_GEMINI", "openai/gemini-2.5-flash")
WHYWHY_MODEL_CHATGPT = os.getenv("WHYWHY_MODEL_CHATGPT", "gpt-5.4")
WHYWHY_MODEL_CLAUDE = os.getenv("WHYWHY_MODEL_CLAUDE", "anthropic/claude-sonnet-4-5")
WHYWHY_MODEL_COSTS_JSON = os.getenv("WHYWHY_MODEL_COSTS_JSON", "")
COLLECTION_NAME = "iatf_knowledge"

WHYWHY_AGENT_CATALOG = {
    "Gemini": {
        "model": WHYWHY_MODEL_GEMINI,
        "role": "Lead Investigator",
        "focus": "広く原因仮説を出し、事実と仮説を分けて一次案を作る",
    },
    "ChatGPT": {
        "model": WHYWHY_MODEL_CHATGPT,
        "role": "Logic Auditor",
        "focus": "5Why の因果の飛躍、表現の曖昧さ、抜け漏れを監査する",
    },
    "Claude": {
        "model": WHYWHY_MODEL_CLAUDE,
        "role": "Countermeasure Critic",
        "focus": "再発防止、標準化、手順・管理策の不足を厳しく見る",
    },
}

WHYWHY_MODE_PRESETS = {
    "省APIモード": {
        "max_agents": 1,
        "use_web": False,
        "synthesis": False,
        "top_k": 3,
        "description": "最小構成です。AI 1名、Web参照なし、統合作業なしで消費を抑えます。",
    },
    "標準モード": {
        "max_agents": 2,
        "use_web": True,
        "synthesis": True,
        "top_k": 5,
        "description": "既定構成です。AI 2名で討議し、最終合意案を作ります。",
    },
    "深掘りモード": {
        "max_agents": 3,
        "use_web": True,
        "synthesis": True,
        "top_k": 6,
        "description": "AI 3名で論点を広げ、より深い監査まで行います。",
    },
}

for d in [WORK_DIR, INGEST_DIR, WIP_DIR, KINDLE_DIR]:
    os.makedirs(d, exist_ok=True)

st.set_page_config(page_title="Clawstack QA Dashboard", layout="wide")

PDCA_STATUS_PATH = "/work/pdca_lab/status.json"


def load_pdca_status():
    try:
        with open(PDCA_STATUS_PATH, "r", encoding="utf-8") as handle:
            return json.load(handle)
    except Exception:
        return {}


pdca_status = load_pdca_status()
with st.sidebar:
    st.markdown("### PDCA Lab")
    if pdca_status:
        st.caption(f"Latest run: {pdca_status.get('latest_run_status', '--')}")
        st.caption(f"Pending review: {pdca_status.get('pending_review_count', 0)}")
        st.caption(f"Production prompt: {pdca_status.get('current_production_prompt_version', '--')}")
    else:
        st.caption("Phase 1 setup pending")
    st.markdown("[Open PDCA Lab](http://localhost:8088/apps/pdca_lab/index.html)")

PLATING_FLAG = os.getenv("ENABLE_PLATING_REFLOW_LAB", "1").lower() not in ("0", "false", "off")
if "/work/scripts" not in sys.path:
    sys.path.append("/work/scripts")

try:
    from plating_quality_analysis import (
        build_initial_case,
        load_defaults as load_plating_defaults,
        load_recent_cases,
        run_analysis as run_plating_analysis,
        save_case as save_plating_case,
        save_uploaded_image as save_plating_uploaded_image,
    )
except Exception:
    build_initial_case = None
    load_plating_defaults = None
    load_recent_cases = None
    run_plating_analysis = None
    save_plating_case = None
    save_plating_uploaded_image = None

# --- UTILS (No external deps) ---

def get_embedding(text):
    try:
        url = f"{OLLAMA_URL}/api/embeddings"
        data = {"model": EMBED_MODEL, "prompt": text}
        req = urllib.request.Request(url, data=json.dumps(data).encode('utf-8'))
        with urllib.request.urlopen(req) as response:
            return json.loads(response.read().decode('utf-8'))['embedding']
    except Exception as e:
        return []

def search_qdrant(vector, limit=3):
    try:
        url = f"{QDRANT_URL}/collections/{COLLECTION_NAME}/points/search"
        data = {"vector": vector, "limit": limit, "with_payload": True}
        req = urllib.request.Request(url, data=json.dumps(data).encode('utf-8'), headers={'Content-Type': 'application/json'})
        with urllib.request.urlopen(req) as response:
            results = json.loads(response.read().decode('utf-8')).get('result', [])
            return "\n\n".join([f"[{r['payload'].get('source', '?')}] {r['payload'].get('text', '')}" for r in results])
    except Exception:
        return "" 

def search_web(query, limit=3):
    try:
        params = urllib.parse.urlencode({
            "q": query,
            "format": "json",
            "engines": "google,bing,duckduckgo",
        })
        with urllib.request.urlopen(f"{SEARXNG_URL}/search?{params}") as response:
            payload = json.loads(response.read().decode("utf-8"))
        results = payload.get("results", [])[:limit]
        return "\n\n".join(
            [
                f"[WEB:{item.get('title', '?')}] {item.get('content', '')}\nURL: {item.get('url', '')}"
                for item in results
            ]
        )
    except Exception:
        return ""

def merge_reference_context(*sections):
    blocks = []
    for title, body in sections:
        if body:
            blocks.append(f"{title}:\n{body}")
    return "\n\n".join(blocks)

def load_whywhy_cost_table():
    if not WHYWHY_MODEL_COSTS_JSON.strip():
        return {}
    try:
        data = json.loads(WHYWHY_MODEL_COSTS_JSON)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}

WHYWHY_MODEL_COSTS = load_whywhy_cost_table()

def estimate_usage_cost(model_name, usage, explicit_cost=None):
    if explicit_cost is not None:
        try:
            return float(explicit_cost)
        except Exception:
            pass
    pricing = WHYWHY_MODEL_COSTS.get(model_name, {})
    if not pricing:
        return None
    prompt_tokens = usage.get("prompt_tokens", 0)
    completion_tokens = usage.get("completion_tokens", 0)
    input_rate = float(pricing.get("input_per_1m_usd", 0))
    output_rate = float(pricing.get("output_per_1m_usd", 0))
    return ((prompt_tokens / 1_000_000) * input_rate) + ((completion_tokens / 1_000_000) * output_rate)

def format_usage_summary(usage_summary):
    cost = usage_summary.get("estimated_cost_usd")
    cost_text = f"${cost:.6f}" if isinstance(cost, (int, float)) else "N/A"
    return (
        f"Prompt {usage_summary.get('prompt_tokens', 0):,} / "
        f"Completion {usage_summary.get('completion_tokens', 0):,} / "
        f"Total {usage_summary.get('total_tokens', 0):,} tokens | "
        f"Estimated cost {cost_text}"
    )

def ask_ai_local(prompt, context_text=""):
    system_prompt = "You are a Quality Assurance Expert."
    if context_text:
        system_prompt += f"\n\nREFERENCE DOCUMENTS:\n{context_text}\n\nUse these references to answer."
    try:
        data = {"model": GEN_MODEL, "prompt": f"{system_prompt}\n\nTask: {prompt}", "stream": False}
        req = urllib.request.Request(f"{OLLAMA_URL}/api/generate", data=json.dumps(data).encode('utf-8'))
        with urllib.request.urlopen(req) as response:
            return json.loads(response.read().decode('utf-8')).get('response', '').strip()
    except Exception as e:
        return f"笞・・AI Offline: {e}"

def ask_ai_deep(prompt, context_text="", model_name="", system_prompt_override=""):
    selected_model = model_name or FMEA_DEEP_MODEL
    if not selected_model:
        return "Deep AI model is not configured."
    system_prompt = system_prompt_override or (
        "You are a senior PFMEA and quality engineering advisor. "
        "Think carefully, prefer concrete manufacturing risk reasoning, and clearly separate assumptions from evidence."
    )
    if context_text:
        system_prompt += f"\n\nREFERENCE DOCUMENTS:\n{context_text}\n\nUse internal references first and use web references only as supplemental evidence."
    body = {
        "model": selected_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ],
        "temperature": 0.2,
    }
    headers = {"Content-Type": "application/json"}
    if FMEA_DEEP_API_KEY:
        headers["Authorization"] = f"Bearer {FMEA_DEEP_API_KEY}"
    req = urllib.request.Request(
        f"{LITELLM_URL}/chat/completions",
        data=json.dumps(body).encode("utf-8"),
        headers=headers,
    )
    with urllib.request.urlopen(req, timeout=90) as response:
        payload = json.loads(response.read().decode("utf-8"))
    choices = payload.get("choices", [])
    if not choices:
        return "Deep AI returned no choices."
    message = choices[0].get("message", {})
    return (message.get("content") or "").strip()


def safe_stem(name: str) -> str:
    base = os.path.splitext(os.path.basename(name))[0]
    normalized = re.sub(r"[^A-Za-z0-9._-]+", "_", base).strip("._-")
    return normalized or "model"


def extract_gdt_pdf_review_context(pdf_path: str, max_pages: int = 6, max_preview_chars: int = 2200) -> dict:
    drawing_info = {
        "drawing_name": os.path.basename(pdf_path),
        "page_count": 0,
        "text_extract_method": "pypdf",
        "text_preview": "",
        "candidate_requirements": [],
        "review_notes": [],
    }
    try:
        reader = pypdf.PdfReader(pdf_path)
        drawing_info["page_count"] = len(reader.pages)
        preview_parts = []
        candidate_lines = []
        keyword_pattern = re.compile(
            r"(datum|profile|position|perpendicular|parallel|flatness|straightness|runout|"
            r"円筒度|真円度|真直度|平面度|直角度|平行度|位置度|同軸度|振れ|データム|幾何公差|基準)",
            re.IGNORECASE,
        )
        for page_index, page in enumerate(reader.pages[:max_pages], start=1):
            page_text = (page.extract_text() or "").replace("\x00", " ")
            compact = re.sub(r"\s+", " ", page_text).strip()
            if not compact:
                continue
            preview_parts.append(f"[Page {page_index}] {compact[:700]}")
            for raw_line in re.split(r"[\r\n]+", page_text):
                line = re.sub(r"\s+", " ", raw_line).strip()
                if len(line) < 4:
                    continue
                if keyword_pattern.search(line):
                    candidate_lines.append(line)
        deduped_candidates = []
        seen = set()
        for line in candidate_lines:
            key = line.lower()
            if key in seen:
                continue
            seen.add(key)
            deduped_candidates.append(line[:180])
            if len(deduped_candidates) >= 12:
                break
        preview_text = "\n\n".join(preview_parts)
        drawing_info["text_preview"] = preview_text[:max_preview_chars]
        drawing_info["candidate_requirements"] = deduped_candidates
        if not deduped_candidates:
            drawing_info["review_notes"].append(
                "No obvious GD&T/datum keywords were extracted automatically; manual drawing review is still required."
            )
    except Exception as exc:
        drawing_info["review_notes"].append(f"PDF extraction failed: {exc}")
    return drawing_info


def build_gdt_requirement_mapping_rows(candidate_requirements):
    rows = []
    for idx, requirement in enumerate(candidate_requirements or [], start=1):
        rows.append({
            "requirement_id": f"REQ-{idx:02d}",
            "requirement_text": requirement,
            "candidate_face_ids": "",
            "candidate_axis_ids": "",
            "chosen_target": "",
            "status": "pending",
            "review_note": "",
        })
    return rows


def summarize_gdt_requirement_mappings(requirement_mappings: list[dict] | None) -> dict:
    rows = requirement_mappings or []
    chosen_count = 0
    pending_count = 0
    rejected_count = 0
    candidate_listed_count = 0
    for row in rows:
        status = (row.get("status") or "").strip().lower()
        chosen_target = (row.get("chosen_target") or "").strip()
        if chosen_target or status == "chosen":
            chosen_count += 1
        elif status == "rejected":
            rejected_count += 1
        elif status == "candidate_listed":
            candidate_listed_count += 1
        else:
            pending_count += 1
    return {
        "total_requirements": len(rows),
        "chosen_count": chosen_count,
        "candidate_listed_count": candidate_listed_count,
        "pending_count": pending_count,
        "rejected_count": rejected_count,
    }


def build_gdt_review_manifest(
    model_name: str,
    html_profile: str,
    drawing_info: dict | None = None,
    requirement_mappings: list[dict] | None = None,
) -> dict:
    drawing_info = drawing_info or {}
    drawing_name = drawing_info.get("drawing_name")
    candidate_requirements = drawing_info.get("candidate_requirements") or []
    requirement_mappings = requirement_mappings or build_gdt_requirement_mapping_rows(candidate_requirements)
    mapping_summary = summarize_gdt_requirement_mappings(requirement_mappings)
    unresolved_points = [
        "Chosen face ids and exact GD&T targets still need engineer review.",
    ]
    if drawing_name:
        unresolved_points.append(
            "Drawing PDF was attached and preview-extracted, but 2D callouts still need mapping to 3D face/axis ids."
        )
    else:
        unresolved_points.append("2D drawing/PDF was not attached in this converter run.")
    return {
        "job_meta": {
            "job_id": f"gdt-review-{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "model_name": model_name,
            "purpose": "GD&T 2D/3D alignment review",
            "artifact_version": "gdt_review_v1",
            "timestamp": datetime.datetime.now(datetime.timezone.utc).isoformat(),
        },
        "conversion": {
            "html_profile": html_profile,
            "viewer_type": "interactive_html",
            "mode": "gdt_review_bundle",
        },
        "drawing": {
            "attached": bool(drawing_name),
            "drawing_name": drawing_name or "",
            "page_count": drawing_info.get("page_count", 0),
            "text_extract_method": drawing_info.get("text_extract_method", ""),
            "candidate_requirement_count": len(candidate_requirements),
            "candidate_requirements": candidate_requirements,
            "text_preview": drawing_info.get("text_preview", ""),
            "review_notes": drawing_info.get("review_notes", []),
        },
        "claim_gate": {
            "render_success": False,
            "placement_verified": False,
            "requires_front_validation": True,
            "requires_side_validation": True,
            "requires_face_or_axis_ids": True,
        },
        "required_outputs": [
            "requirement_list",
            "candidate_face_ids",
            "chosen_face_ids_or_axis_ids",
            "front_validation",
            "side_validation",
            "unresolved_points",
        ],
        "validation": {
            "front": {"status": "pending", "notes": ""},
            "side": {"status": "pending", "notes": ""},
            "top": {"status": "optional", "notes": ""},
        },
        "review_mapping": {
            "requirement_rows": requirement_mappings,
            "mapping_complete": mapping_summary["chosen_count"] > 0 and mapping_summary["pending_count"] == 0,
            "summary": mapping_summary,
        },
        "status": {
            "render_success": True,
            "placement_verified": False,
            "unresolved_points": unresolved_points,
        },
    }


def build_gdt_checklist_markdown(model_name: str, drawing_name: str = "") -> str:
    lines = [
        f"# GD&T 2D/3D Alignment Checklist",
        "",
        f"Model: {model_name}",
        f"Drawing: {drawing_name or 'Not attached'}",
        f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        "",
        "## Drawing Read",
        "- [ ] datum / section / detail symbols were separated correctly",
        "- [ ] drawing target face or axis is identified from the 2D document",
        "- [ ] theoretical geometry and physical faces are not mixed",
        "",
        "## 3D Mapping",
        "- [ ] candidate face ids were listed",
        "- [ ] chosen face ids / axis ids were recorded",
        "- [ ] exact layer and debug layer are separated",
        "",
        "## Validation",
        "- [ ] front view is acceptable",
        "- [ ] side view is acceptable",
        "- [ ] top or local crop was checked if ambiguity remains",
        "- [ ] explanation text matches the rendered geometry",
        "",
        "## Before Claiming Success",
        "- [ ] unresolved items are explicitly listed",
        "- [ ] no 'STEP-face based' claim is made without face ids",
        "- [ ] screenshot evidence was reviewed",
        "",
    ]
    return "\n".join(lines)

def ask_ai_deep_with_meta(prompt, context_text="", model_name="", system_prompt_override=""):
    selected_model = model_name or FMEA_DEEP_MODEL
    if not selected_model:
        return {
            "content": "Deep AI model is not configured.",
            "model": selected_model,
            "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
            "estimated_cost_usd": None,
            "error": "model_not_configured",
        }
    system_prompt = system_prompt_override or (
        "You are a senior PFMEA and quality engineering advisor. "
        "Think carefully, prefer concrete manufacturing risk reasoning, and clearly separate assumptions from evidence."
    )
    if context_text:
        system_prompt += f"\n\nREFERENCE DOCUMENTS:\n{context_text}\n\nUse internal references first and use web references only as supplemental evidence."
    body = {
        "model": selected_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ],
        "temperature": 0.2,
    }
    headers = {"Content-Type": "application/json"}
    if FMEA_DEEP_API_KEY:
        headers["Authorization"] = f"Bearer {FMEA_DEEP_API_KEY}"
    req = urllib.request.Request(
        f"{LITELLM_URL}/chat/completions",
        data=json.dumps(body).encode("utf-8"),
        headers=headers,
    )
    with urllib.request.urlopen(req, timeout=90) as response:
        payload = json.loads(response.read().decode("utf-8"))
    choices = payload.get("choices", [])
    message = choices[0].get("message", {}) if choices else {}
    usage = payload.get("usage", {}) or {}
    usage_summary = {
        "prompt_tokens": int(usage.get("prompt_tokens", 0) or 0),
        "completion_tokens": int(usage.get("completion_tokens", 0) or 0),
        "total_tokens": int(usage.get("total_tokens", 0) or 0),
    }
    estimated_cost = estimate_usage_cost(
        selected_model,
        usage_summary,
        explicit_cost=payload.get("response_cost", usage.get("cost")),
    )
    return {
        "content": (message.get("content") or "").strip(),
        "model": selected_model,
        "usage": usage_summary,
        "estimated_cost_usd": estimated_cost,
        "error": None,
    }

def ask_ai(prompt, context_text="", mode="local", deep_model_name="", system_prompt_override=""):
    if mode == "deep":
        try:
            deep_reply = ask_ai_deep(
                prompt,
                context_text=context_text,
                model_name=deep_model_name,
                system_prompt_override=system_prompt_override,
            )
            if deep_reply:
                return deep_reply
        except Exception as e:
            fallback = ask_ai_local(prompt, context_text=context_text)
            return f"[Deep AI unavailable: {e}]\n\n{fallback}"
    return ask_ai_local(prompt, context_text=context_text)

def build_whywhy_context(problem, whys, use_internal_docs=True, use_web_docs=True, rag_limit=5):
    query_parts = [problem, *[item for item in whys if item]]
    query = " | ".join([part for part in query_parts if part]).strip()
    if not query:
        return "", "", "", ""
    rag_context = ""
    if use_internal_docs:
        vec = get_embedding(query)
        rag_context = search_qdrant(vec, limit=rag_limit) if vec else ""
    web_context = ""
    if use_web_docs:
        web_context = search_web(f'5 why root cause analysis manufacturing quality {problem}', limit=3)
    context = merge_reference_context(
        ("MITSUI / INTERNAL QUALITY KNOWLEDGE", rag_context),
        ("PUBLIC WEB KNOWLEDGE", web_context),
    )
    return query, rag_context, web_context, context

def run_whywhy_agents(problem, whys, selected_agent_names, context_text, synthesis_enabled=True):
    steps_text = "\n".join([f"{index}. {item}" for index, item in enumerate(whys, start=1) if item]) or "(No why steps entered)"
    role_plan = {
        1: "Lead Investigator",
        2: "Logic Auditor",
        3: "Countermeasure Critic",
    }
    agent_outputs = []
    usage_rollup = {
        "prompt_tokens": 0,
        "completion_tokens": 0,
        "total_tokens": 0,
        "estimated_cost_usd": 0.0,
        "cost_known": False,
        "calls": 0,
    }
    for index, agent_name in enumerate(selected_agent_names, start=1):
        agent = WHYWHY_AGENT_CATALOG.get(agent_name, {})
        assigned_role = role_plan.get(index, agent.get("role", "Reviewer"))
        system_prompt = (
            f"You are participating in a 5-Why quality review as {agent_name}. "
            f"Your assigned role is {assigned_role}. "
            f"Your focus is: {agent.get('focus', 'analyze carefully')}. "
            "Be practical for manufacturing quality management, separate evidence from assumptions, "
            "and point out weak causal links or missing verification."
        )
        prompt = f"""
Problem:
{problem}

Current 5-Why draft:
{steps_text}

Task:
1. Evaluate whether the causal chain is logically connected.
2. Point out weak or missing links.
3. Suggest improved why statements if needed.
4. Suggest what evidence or standard/procedure should be checked next.
5. Run a backward check: start from the assumed deepest cause and verify whether the chain can realistically return to the original problem.
6. If the backward check fails, point out the exact why-step where the reverse logic breaks.
7. Keep the response concise and structured.
"""
        reply_meta = ask_ai_deep_with_meta(
            prompt,
            context_text=context_text,
            system_prompt_override=system_prompt,
            model_name=agent.get("model", ""),
        )
        usage = reply_meta.get("usage", {})
        usage_rollup["prompt_tokens"] += usage.get("prompt_tokens", 0)
        usage_rollup["completion_tokens"] += usage.get("completion_tokens", 0)
        usage_rollup["total_tokens"] += usage.get("total_tokens", 0)
        usage_rollup["calls"] += 1
        if isinstance(reply_meta.get("estimated_cost_usd"), (int, float)):
            usage_rollup["estimated_cost_usd"] += float(reply_meta["estimated_cost_usd"])
            usage_rollup["cost_known"] = True
        agent_outputs.append({
            "agent_name": agent_name,
            "assigned_role": assigned_role,
            "model": reply_meta.get("model", agent.get("model", "")),
            "focus": agent.get("focus", ""),
            "reply": reply_meta.get("content", ""),
            "usage": usage,
            "estimated_cost_usd": reply_meta.get("estimated_cost_usd"),
        })

    if not synthesis_enabled:
        summary = agent_outputs[0]["reply"] if agent_outputs else "No agent output."
        usage_rollup["estimated_cost_usd"] = usage_rollup["estimated_cost_usd"] if usage_rollup["cost_known"] else None
        return agent_outputs, summary, usage_rollup

    synthesis_prompt = f"""
Problem:
{problem}

Current 5-Why draft:
{steps_text}

Agent reviews:
{chr(10).join([f"[{item['agent_name']} - {item['assigned_role']}]{chr(10)}{item['reply']}" for item in agent_outputs])}

Task:
Create a final integrated 5-Why review with these sections:
1. Consensus Summary
2. Rewritten 5-Why Draft
3. Backward Validation
4. Disagreements or open points
5. Evidence to confirm next
6. Recommended containment / corrective action direction
"""
    moderator_name = selected_agent_names[0] if selected_agent_names else "Gemini"
    moderator = WHYWHY_AGENT_CATALOG.get(moderator_name, {})
    synthesis_meta = ask_ai_deep_with_meta(
        synthesis_prompt,
        context_text=context_text,
        system_prompt_override=(
            "You are the moderator of a multi-agent 5-Why quality review. "
            "Synthesize competing viewpoints fairly, prefer internal standards when available, "
            "and make the final output practical for a manufacturing quality team."
        ),
        model_name=moderator.get("model", ""),
    )
    synthesis_usage = synthesis_meta.get("usage", {})
    usage_rollup["prompt_tokens"] += synthesis_usage.get("prompt_tokens", 0)
    usage_rollup["completion_tokens"] += synthesis_usage.get("completion_tokens", 0)
    usage_rollup["total_tokens"] += synthesis_usage.get("total_tokens", 0)
    usage_rollup["calls"] += 1
    if isinstance(synthesis_meta.get("estimated_cost_usd"), (int, float)):
        usage_rollup["estimated_cost_usd"] += float(synthesis_meta["estimated_cost_usd"])
        usage_rollup["cost_known"] = True
    usage_rollup["estimated_cost_usd"] = usage_rollup["estimated_cost_usd"] if usage_rollup["cost_known"] else None
    return agent_outputs, synthesis_meta.get("content", ""), usage_rollup

def build_standard_review_context(query, web_query, use_internal_docs=True, use_web_docs=True, rag_limit=5):
    query = (query or "").strip()
    if not query:
        return "", "", "", ""
    rag_context = ""
    if use_internal_docs:
        vec = get_embedding(query)
        rag_context = search_qdrant(vec, limit=rag_limit) if vec else ""
    web_context = ""
    if use_web_docs:
        web_context = search_web(web_query, limit=3)
    context = merge_reference_context(
        ("MITSUI / INTERNAL QUALITY KNOWLEDGE", rag_context),
        ("PUBLIC WEB KNOWLEDGE", web_context),
    )
    return query, rag_context, web_context, context

def run_multi_agent_quality_review(
    subject_title,
    subject_body,
    selected_agent_names,
    context_text,
    review_task_text,
    synthesis_task_text,
    synthesis_enabled=True,
):
    role_plan = {
        1: "Lead Investigator",
        2: "Logic Auditor",
        3: "Countermeasure Critic",
    }
    agent_outputs = []
    usage_rollup = {
        "prompt_tokens": 0,
        "completion_tokens": 0,
        "total_tokens": 0,
        "estimated_cost_usd": 0.0,
        "cost_known": False,
        "calls": 0,
    }
    for index, agent_name in enumerate(selected_agent_names, start=1):
        agent = WHYWHY_AGENT_CATALOG.get(agent_name, {})
        assigned_role = role_plan.get(index, agent.get("role", "Reviewer"))
        system_prompt = (
            f"You are participating in a multi-AI manufacturing quality review as {agent_name}. "
            f"Your assigned role is {assigned_role}. "
            f"Your focus is: {agent.get('focus', 'analyze carefully')}. "
            "Use internal standards first when available, separate evidence from assumptions, "
            "and make the review practical for a quality engineering team."
        )
        prompt = f"""
{subject_title}:
{subject_body}

Task:
{review_task_text}
"""
        reply_meta = ask_ai_deep_with_meta(
            prompt,
            context_text=context_text,
            system_prompt_override=system_prompt,
            model_name=agent.get("model", ""),
        )
        usage = reply_meta.get("usage", {})
        usage_rollup["prompt_tokens"] += usage.get("prompt_tokens", 0)
        usage_rollup["completion_tokens"] += usage.get("completion_tokens", 0)
        usage_rollup["total_tokens"] += usage.get("total_tokens", 0)
        usage_rollup["calls"] += 1
        if isinstance(reply_meta.get("estimated_cost_usd"), (int, float)):
            usage_rollup["estimated_cost_usd"] += float(reply_meta["estimated_cost_usd"])
            usage_rollup["cost_known"] = True
        agent_outputs.append({
            "agent_name": agent_name,
            "assigned_role": assigned_role,
            "model": reply_meta.get("model", agent.get("model", "")),
            "focus": agent.get("focus", ""),
            "reply": reply_meta.get("content", ""),
            "usage": usage,
            "estimated_cost_usd": reply_meta.get("estimated_cost_usd"),
        })

    if not synthesis_enabled:
        summary = agent_outputs[0]["reply"] if agent_outputs else "No agent output."
        usage_rollup["estimated_cost_usd"] = usage_rollup["estimated_cost_usd"] if usage_rollup["cost_known"] else None
        return agent_outputs, summary, usage_rollup

    moderator_name = selected_agent_names[0] if selected_agent_names else "Gemini"
    moderator = WHYWHY_AGENT_CATALOG.get(moderator_name, {})
    synthesis_prompt = f"""
{subject_title}:
{subject_body}

Agent reviews:
{chr(10).join([f"[{item['agent_name']} - {item['assigned_role']}]{chr(10)}{item['reply']}" for item in agent_outputs])}

Task:
{synthesis_task_text}
"""
    synthesis_meta = ask_ai_deep_with_meta(
        synthesis_prompt,
        context_text=context_text,
        system_prompt_override=(
            "You are the moderator of a multi-agent quality engineering review. "
            "Synthesize competing viewpoints fairly, prefer internal standards when available, "
            "and produce a practical conclusion for a manufacturing team."
        ),
        model_name=moderator.get("model", ""),
    )
    synthesis_usage = synthesis_meta.get("usage", {})
    usage_rollup["prompt_tokens"] += synthesis_usage.get("prompt_tokens", 0)
    usage_rollup["completion_tokens"] += synthesis_usage.get("completion_tokens", 0)
    usage_rollup["total_tokens"] += synthesis_usage.get("total_tokens", 0)
    usage_rollup["calls"] += 1
    if isinstance(synthesis_meta.get("estimated_cost_usd"), (int, float)):
        usage_rollup["estimated_cost_usd"] += float(synthesis_meta["estimated_cost_usd"])
        usage_rollup["cost_known"] = True
    usage_rollup["estimated_cost_usd"] = usage_rollup["estimated_cost_usd"] if usage_rollup["cost_known"] else None
    return agent_outputs, synthesis_meta.get("content", ""), usage_rollup

def default_pfmea_rows(process_step):
    return [{
        "Process Step": process_step,
        "Process Function": "",
        "Requirement": "",
        "Potential Failure Mode": "",
        "Potential Effect": "",
        "Severity": 0,
        "Potential Cause": "",
        "Occurrence": 0,
        "Current Prevention Control": "",
        "Current Detection Control": "",
        "Detection": 0,
        "RPN": 0,
        "Recommended Action": "",
        "Responsibility": "",
        "Due Date": "",
        "Action Status": "Open",
    }]

def normalize_pfmea_dataframe(df, process_step):
    expected_columns = [
        "Process Step",
        "Process Function",
        "Requirement",
        "Potential Failure Mode",
        "Potential Effect",
        "Severity",
        "Potential Cause",
        "Occurrence",
        "Current Prevention Control",
        "Current Detection Control",
        "Detection",
        "RPN",
        "Recommended Action",
        "Responsibility",
        "Due Date",
        "Action Status",
    ]
    if df is None or df.empty:
        return pd.DataFrame(default_pfmea_rows(process_step))
    normalized = df.copy()
    for column in expected_columns:
        if column not in normalized.columns:
            normalized[column] = ""
    numeric_columns = ["Severity", "Occurrence", "Detection"]
    for column in numeric_columns:
        normalized[column] = pd.to_numeric(normalized[column], errors="coerce").fillna(0).astype(int).clip(lower=0, upper=10)
    normalized["Process Step"] = normalized["Process Step"].replace("", process_step).fillna(process_step)
    normalized["RPN"] = normalized["Severity"] * normalized["Occurrence"] * normalized["Detection"]
    normalized["Action Status"] = normalized["Action Status"].replace("", "Open").fillna("Open")
    return normalized[expected_columns]

def extract_text_immediate(filepath):
    """Refactored extraction logic for immediate use"""
    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(filepath)
            for p in reader.pages: text += p.extract_text() + "\n"
        elif ext in [".docx", ".doc"]:
            doc = docx.Document(filepath)
            for p in doc.paragraphs: text += p.text + "\n"
        elif ext in [".xlsx", ".xls"]:
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for s in wb.sheetnames:
                for r in wb[s].iter_rows(values_only=True):
                    text += " ".join([str(c) for c in r if c]) + "\n"
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f: text = f.read()
    except Exception as e:
        return f"[Error extracting {ext}: {e}]"
    return text

def save_uploaded_file(uploaded_file, target_folder):
    try:
        file_path = os.path.join(target_folder, uploaded_file.name)
        if os.path.exists(file_path):
            return True, file_path, "Exists"
        
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        return True, file_path, "Saved"
    except Exception as e:
        return False, str(e), "Error"


def render_plating_quality_page():
    st.header("Plating Quality Analysis")
    st.caption("Thermal FEM plus reduced-order plating/reflow indicators.")

    required = [
        build_initial_case,
        load_plating_defaults,
        load_recent_cases,
        run_plating_analysis,
        save_plating_case,
        save_plating_uploaded_image,
    ]
    if not all(required):
        st.error("plating_quality_analysis.py could not be loaded. Check /work/scripts.")
        return

    defaults = load_plating_defaults()
    if "plating_case" not in st.session_state:
        st.session_state.plating_case = build_initial_case(defaults)
    if "plating_results" not in st.session_state:
        st.session_state.plating_results = {}
    if "plating_images" not in st.session_state:
        st.session_state.plating_images = []

    case = st.session_state.plating_case
    results = st.session_state.plating_results

    top1, top2, top3, top4 = st.columns(4)
    top1.metric("Project", case["project"]["project_name"])
    top2.metric("Stackup", case["project"]["stackup"])
    top3.metric("Status", case["project"]["status"])
    top4.metric("Mode", case["analysis"]["analysis_mode"])

    tab_input, tab_images, tab_result, tab_history = st.tabs(["Inputs", "Assets", "Results", "History"])

    with tab_input:
        left, right = st.columns(2)
        with left:
            st.subheader("Project")
            case["project"]["project_name"] = st.text_input("Project name", case["project"]["project_name"])
            case["project"]["part_number"] = st.text_input("Part number", case["project"]["part_number"])
            case["project"]["revision"] = st.text_input("Revision", case["project"]["revision"])
            case["project"]["material_system"] = st.text_input("Material system", case["project"]["material_system"])
            case["project"]["substrate"] = st.text_input("Substrate", case["project"]["substrate"])
            case["project"]["stackup"] = st.text_input("Stackup", case["project"]["stackup"])
            status_options = ["draft", "ready", "running", "done", "error"]
            status_value = case["project"]["status"] if case["project"]["status"] in status_options else "draft"
            case["project"]["status"] = st.selectbox("Status", status_options, index=status_options.index(status_value))

            st.subheader("Specification")
            case["spec"]["substrate_type"] = st.text_input("Substrate type", case["spec"]["substrate_type"])
            case["spec"]["substrate_grade"] = st.text_input("Substrate grade", case["spec"]["substrate_grade"])
            case["spec"]["substrate_thickness_um"] = st.number_input("Substrate thickness (um)", value=float(case["spec"]["substrate_thickness_um"]), step=10.0)
            case["spec"]["ep_layer_enabled"] = st.checkbox("EP layer enabled", value=bool(case["spec"]["ep_layer_enabled"]))
            case["spec"]["ni_thickness_target_um"] = st.number_input("Ni target thickness (um)", value=float(case["spec"]["ni_thickness_target_um"]), step=0.05, format="%.3f")
            case["spec"]["sn_thickness_target_um"] = st.number_input("Sn target thickness (um)", value=float(case["spec"]["sn_thickness_target_um"]), step=0.05, format="%.3f")
            case["spec"]["ni_sn_max_um"] = st.number_input("Ni-Sn max (um)", value=float(case["spec"]["ni_sn_max_um"]), step=0.1)
            case["spec"]["initial_imc_thickness_um"] = st.number_input("Initial IMC thickness (um)", value=float(case["spec"]["initial_imc_thickness_um"]), step=0.01, format="%.3f")
            case["spec"]["surface_roughness_ra_um"] = st.number_input("Surface roughness Ra (um)", value=float(case["spec"]["surface_roughness_ra_um"]), step=0.01, format="%.3f")
            case["spec"]["plating_side_mode"] = st.selectbox("Plating side mode", ["single", "double"], index=0 if case["spec"]["plating_side_mode"] == "single" else 1)
            case["spec"]["note_spec"] = st.text_area("Notes", case["spec"].get("note_spec", ""), height=80)

        with right:
            st.subheader("Plating")
            case["plating"]["plating_line_type"] = st.text_input("Line type", case["plating"]["plating_line_type"])
            case["plating"]["plating_machine_name"] = st.text_input("Machine name", case["plating"]["plating_machine_name"])
            case["plating"]["plating_bath_name"] = st.text_input("Bath name", case["plating"]["plating_bath_name"])
            case["plating"]["plating_current_density_adm2"] = st.number_input("Current density (A/dm2)", min_value=1.0, max_value=30.0, value=float(case["plating"]["plating_current_density_adm2"]), step=0.1)
            case["plating"]["plating_current_mode"] = st.text_input("Current mode", case["plating"]["plating_current_mode"])
            case["plating"]["plating_line_speed_m_min"] = st.number_input("Line speed (m/min)", value=float(case["plating"]["plating_line_speed_m_min"]), step=0.1)
            case["plating"]["plating_bath_temp_c"] = st.number_input("Bath temp (C)", value=float(case["plating"]["plating_bath_temp_c"]), step=1.0)
            case["plating"]["plating_time_sec"] = st.number_input("Plating time (sec)", value=float(case["plating"]["plating_time_sec"]), step=1.0)
            case["plating"]["agitation_mode"] = st.text_input("Agitation mode", case["plating"]["agitation_mode"])
            case["plating"]["anode_type"] = st.text_input("Anode type", case["plating"]["anode_type"])
            case["plating"]["xrf_measurement_enabled"] = st.checkbox("XRF enabled", value=bool(case["plating"]["xrf_measurement_enabled"]))
            case["plating"]["xrf_points_count"] = st.number_input("XRF points", value=int(case["plating"]["xrf_points_count"]), step=1)
            case["plating"]["thickness_uniformity_index"] = st.number_input("Uniformity index", value=float(case["plating"]["thickness_uniformity_index"]), step=0.1)
            case["plating"]["surface_orientation_note"] = st.text_area("Orientation note", case["plating"]["surface_orientation_note"], height=80)

            st.subheader("Reflow")
            case["reflow"]["reflow_machine_type"] = st.text_input("Machine type", case["reflow"]["reflow_machine_type"])
            case["reflow"]["reflow_machine_name"] = st.text_input("Machine name", case["reflow"]["reflow_machine_name"])
            atm_options = ["air", "nitrogen", "vacuum"]
            atm = case["reflow"]["atmosphere_type"] if case["reflow"]["atmosphere_type"] in atm_options else "air"
            case["reflow"]["atmosphere_type"] = st.selectbox("Atmosphere", atm_options, index=atm_options.index(atm))
            case["reflow"]["o2_ppm"] = st.number_input("O2 (ppm)", value=float(case["reflow"]["o2_ppm"]), step=100.0)
            case["reflow"]["zones_count"] = st.number_input("Zone count", value=int(case["reflow"]["zones_count"]), step=1)
            case["reflow"]["conveyor_speed_mm_min"] = st.number_input("Conveyor speed (mm/min)", value=float(case["reflow"]["conveyor_speed_mm_min"]), step=10.0)
            case["reflow"]["board_or_carrier_type"] = st.text_input("Board/carrier", case["reflow"]["board_or_carrier_type"])
            case["reflow"]["flux_or_residue_condition"] = st.text_input("Flux/residue condition", case["reflow"]["flux_or_residue_condition"])
            case["reflow"]["reflow_repeat_count"] = st.number_input("Reflow repeat count", value=int(case["reflow"]["reflow_repeat_count"]), step=1)
            case["reflow"]["cooling_mode"] = st.text_input("Cooling mode", case["reflow"]["cooling_mode"])

        low_left, low_right = st.columns(2)
        with low_left:
            st.subheader("Temperature Profile")
            case["profile"]["profile_template_name"] = st.text_input("Template", case["profile"]["profile_template_name"])
            case["profile"]["start_temp_c"] = st.number_input("Start temp (C)", value=float(case["profile"]["start_temp_c"]), step=1.0)
            case["profile"]["preheat_target_c"] = st.number_input("Preheat target (C)", value=float(case["profile"]["preheat_target_c"]), step=1.0)
            case["profile"]["preheat_time_sec"] = st.number_input("Preheat time (sec)", value=float(case["profile"]["preheat_time_sec"]), step=5.0)
            case["profile"]["soak_min_c"] = st.number_input("Soak min (C)", value=float(case["profile"]["soak_min_c"]), step=1.0)
            case["profile"]["soak_max_c"] = st.number_input("Soak max (C)", value=float(case["profile"]["soak_max_c"]), step=1.0)
            case["profile"]["soak_time_sec"] = st.number_input("Soak time (sec)", value=float(case["profile"]["soak_time_sec"]), step=5.0)
            case["profile"]["ramp_to_peak_sec"] = st.number_input("Ramp to peak (sec)", value=float(case["profile"]["ramp_to_peak_sec"]), step=5.0)
            case["profile"]["peak_temp_c"] = st.number_input("Peak temp (C)", value=float(case["profile"]["peak_temp_c"]), step=1.0)
            case["profile"]["tal_over_liquidus_sec"] = st.number_input("TAL (sec)", value=float(case["profile"]["tal_over_liquidus_sec"]), step=1.0)
            case["profile"]["liquidus_temp_c"] = st.number_input("Liquidus temp (C)", value=float(case["profile"]["liquidus_temp_c"]), step=1.0)
            case["profile"]["cool_to_temp_c"] = st.number_input("Cool to (C)", value=float(case["profile"]["cool_to_temp_c"]), step=1.0)
            case["profile"]["cool_time_sec"] = st.number_input("Cool time (sec)", value=float(case["profile"]["cool_time_sec"]), step=5.0)
            case["profile"]["ramp_rate_c_per_sec"] = st.number_input("Ramp rate (C/sec)", value=float(case["profile"]["ramp_rate_c_per_sec"]), step=0.1)
            case["profile"]["cool_rate_c_per_sec"] = st.number_input("Cool rate (C/sec)", value=float(case["profile"]["cool_rate_c_per_sec"]), step=0.1)

        with low_right:
            st.subheader("Analysis")
            mode_options = ["plating_plus_reflow_coupled", "reflow_only", "plating_only"]
            mode = case["analysis"]["analysis_mode"] if case["analysis"]["analysis_mode"] in mode_options else mode_options[0]
            case["analysis"]["analysis_mode"] = st.selectbox("Analysis mode", mode_options, index=mode_options.index(mode))
            dim_options = ["1D", "2D", "3D"]
            dim = case["analysis"]["model_dimension"] if case["analysis"]["model_dimension"] in dim_options else "2D"
            case["analysis"]["model_dimension"] = st.selectbox("Model dimension", dim_options, index=dim_options.index(dim))
            solver_options = ["scikit-fem", "pycalphad", "elmer", "openfoam"]
            solver = case["analysis"]["solver_backend"] if case["analysis"]["solver_backend"] in solver_options else "scikit-fem"
            case["analysis"]["solver_backend"] = st.selectbox("Solver backend", solver_options, index=solver_options.index(solver))
            case["analysis"]["use_pycalphad"] = st.checkbox("Use pycalphad", value=bool(case["analysis"]["use_pycalphad"]))
            case["analysis"]["use_scikit_fem"] = st.checkbox("Use scikit-fem", value=bool(case["analysis"]["use_scikit_fem"]))
            case["analysis"]["use_fenicsx_if_available"] = st.checkbox("Use FEniCSx if available", value=bool(case["analysis"]["use_fenicsx_if_available"]))
            case["analysis"]["use_openfoam_if_available"] = st.checkbox("Use OpenFOAM if available", value=bool(case["analysis"]["use_openfoam_if_available"]))
            case["analysis"]["use_calculix_if_available"] = st.checkbox("Use CalculiX if available", value=bool(case["analysis"]["use_calculix_if_available"]))
            case["analysis"]["use_elmer_if_available"] = st.checkbox("Use Elmer if available", value=bool(case["analysis"]["use_elmer_if_available"]))
            case["analysis"]["use_paraview_export"] = st.checkbox("Export ParaView artifacts", value=bool(case["analysis"]["use_paraview_export"]))
            case["analysis"]["mesh_size_um"] = st.number_input("Mesh size (um)", value=float(case["analysis"]["mesh_size_um"]), step=0.01, format="%.3f")
            case["analysis"]["time_step_sec"] = st.number_input("Time step (sec)", value=float(case["analysis"]["time_step_sec"]), step=0.1)
            case["analysis"]["total_sim_time_sec"] = st.number_input("Total sim time (sec)", value=float(case["analysis"]["total_sim_time_sec"]), step=10.0)
            case["analysis"]["thermal_coupling_enabled"] = st.checkbox("Thermal coupling", value=bool(case["analysis"]["thermal_coupling_enabled"]))
            case["analysis"]["diffusion_enabled"] = st.checkbox("Diffusion enabled", value=bool(case["analysis"]["diffusion_enabled"]))
            case["analysis"]["imc_growth_enabled"] = st.checkbox("IMC growth enabled", value=bool(case["analysis"]["imc_growth_enabled"]))
            case["analysis"]["void_risk_enabled"] = st.checkbox("Void risk enabled", value=bool(case["analysis"]["void_risk_enabled"]))
            case["analysis"]["adhesion_risk_enabled"] = st.checkbox("Adhesion risk enabled", value=bool(case["analysis"]["adhesion_risk_enabled"]))

        a1, a2, a3 = st.columns([1, 1, 2])
        with a1:
            if st.button("Save case", use_container_width=True):
                case["project"]["updated_at"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
                save_plating_case(case, results, st.session_state.plating_images)
                st.success("Case saved.")
        with a2:
            if st.button("Run analysis", use_container_width=True, type="primary"):
                case["project"]["status"] = "running"
                case["project"]["updated_at"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
                st.session_state.plating_results = run_plating_analysis(case)
                case["project"]["status"] = "done"
                save_plating_case(case, st.session_state.plating_results, st.session_state.plating_images)
                st.success("Analysis completed.")
                st.rerun()
        with a3:
            st.caption("Thermal FEM is solved explicitly. Diffusion, IMC growth, liquid fraction, and crystal order are currently reduced-order proxies.")

    with tab_images:
        st.subheader("Observation Assets")
        upload_col, meta_col = st.columns([1, 1])
        with upload_col:
            uploaded_images = st.file_uploader("Upload SEM / FIB / EDX / XRF / CSV", type=["png", "jpg", "jpeg", "tif", "tiff", "csv"], accept_multiple_files=True, key="plating_image_upload")
        with meta_col:
            location_tag = st.text_input("location_tag", "front-center")
            magnification = st.text_input("magnification", "500x")
            pre_or_post = st.selectbox("pre_or_post_reflow", ["pre", "post"])
            sample_id = st.text_input("sample_id", case["project"]["part_number"])
            image_note = st.text_area("note", "", height=80)

        if uploaded_images and st.button("Save uploaded assets", type="primary"):
            saved = 0
            for item in uploaded_images:
                saved_path = save_plating_uploaded_image(case["project"]["case_id"], item)
                st.session_state.plating_images.append({
                    "file_name": item.name,
                    "saved_path": saved_path,
                    "location_tag": location_tag,
                    "magnification": magnification,
                    "pre_or_post_reflow": pre_or_post,
                    "sample_id": sample_id,
                    "note": image_note,
                })
                saved += 1
            save_plating_case(case, results, st.session_state.plating_images)
            st.success(f"Saved {saved} assets.")

        if st.session_state.plating_images:
            st.dataframe(pd.DataFrame(st.session_state.plating_images), use_container_width=True, hide_index=True)

    with tab_result:
        st.subheader("Results")
        if not results:
            st.info("Run analysis to generate outputs.")
        else:
            fidelity = results.get("model_fidelity", {})
            if fidelity:
                st.info("Thermal field uses FEM. Other fields remain reduced-order for now.")
                st.json(fidelity, expanded=False)

            preds = results["predictions"]
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("Sn remaining (um)", preds["predicted_sn_remaining_um"])
            m2.metric("Ni remaining (um)", preds["predicted_ni_remaining_um"])
            m3.metric("IMC thickness (um)", preds["predicted_imc_thickness_um"])
            m4.metric("Peak stress (MPa)", preds["predicted_peak_stress_mpa"])

            r1, r2, r3, r4 = st.columns(4)
            r1.metric("Void risk", preds["predicted_void_risk_score"], preds["predicted_void_risk_class"])
            r2.metric("Adhesion risk", preds["predicted_adhesion_risk_score"], preds["predicted_adhesion_risk_class"])
            r3.metric("Surface melt risk", preds["predicted_surface_melt_score"], preds["predicted_surface_melt_warning_badge"])
            r4.metric("Wetting score", preds["predicted_wetting_score"])

            derived = results["derived"]
            st.write(f"TAL: **{derived['tal_recomputed_sec']} sec** / Melt score: **{derived['melt_score']}** / Peak: **{derived['peak_temp_c']} C**")

            curve_df = pd.DataFrame(results["profile_curve"])
            st.line_chart(curve_df.rename(columns={"time_sec": "index"}).set_index("index")[["temp_c"]], use_container_width=True)
            st.dataframe(pd.DataFrame([preds]), use_container_width=True, hide_index=True)

            artifacts = results.get("artifacts", {})
            if artifacts:
                st.markdown("---")
                st.subheader("ParaView / VTK Output")
                b1, b2, b3 = st.columns(3)
                with b1:
                    st.caption(f"Output dir: `{artifacts.get('output_dir', '')}`")
                    if artifacts.get("field_vtu") and os.path.exists(artifacts["field_vtu"]):
                        with open(artifacts["field_vtu"], "rb") as handle:
                            st.download_button("Download VTU", handle.read(), file_name=os.path.basename(artifacts["field_vtu"]), use_container_width=True)
                with b2:
                    if artifacts.get("profile_csv") and os.path.exists(artifacts["profile_csv"]):
                        with open(artifacts["profile_csv"], "rb") as handle:
                            st.download_button("Download profile CSV", handle.read(), file_name=os.path.basename(artifacts["profile_csv"]), use_container_width=True)
                with b3:
                    if artifacts.get("summary_json") and os.path.exists(artifacts["summary_json"]):
                        with open(artifacts["summary_json"], "rb") as handle:
                            st.download_button("Download summary JSON", handle.read(), file_name=os.path.basename(artifacts["summary_json"]), use_container_width=True)

                c1, c2 = st.columns(2)
                with c1:
                    if artifacts.get("timeline_pvd") and os.path.exists(artifacts["timeline_pvd"]):
                        with open(artifacts["timeline_pvd"], "rb") as handle:
                            st.download_button("Download PVD timeline", handle.read(), file_name=os.path.basename(artifacts["timeline_pvd"]), use_container_width=True)
                with c2:
                    if artifacts.get("timeline_summary_json") and os.path.exists(artifacts["timeline_summary_json"]):
                        with open(artifacts["timeline_summary_json"], "rb") as handle:
                            st.download_button("Download timeline JSON", handle.read(), file_name=os.path.basename(artifacts["timeline_summary_json"]), use_container_width=True)

                if artifacts.get("snapshot_count"):
                    st.caption(f"Transient VTU snapshots: {artifacts['snapshot_count']}")
                if artifacts.get("paraview_preview_png") and os.path.exists(artifacts["paraview_preview_png"]):
                    st.image(artifacts["paraview_preview_png"], caption="ParaView-compatible preview", use_container_width=True)
                else:
                    st.warning("Preview PNG was not generated, but VTU/PVD output is available.")

    with tab_history:
        st.subheader("Recent cases")
        recent = load_recent_cases(limit=10)
        if not recent:
            st.info("No saved cases found.")
        else:
            history_rows = []
            for item in recent:
                payload = item.get("case", {})
                history_rows.append({
                    "case_id": payload.get("project", {}).get("case_id", ""),
                    "project_name": payload.get("project", {}).get("project_name", ""),
                    "status": payload.get("project", {}).get("status", ""),
                    "updated_at": payload.get("project", {}).get("updated_at", ""),
                    "solver": payload.get("analysis", {}).get("solver_backend", ""),
                })
            st.dataframe(pd.DataFrame(history_rows), use_container_width=True, hide_index=True)


def render_process_monitoring_measurement_page():
    st.header("Process Monitoring & Measurement")
    st.caption("Open the existing IATF process monitoring and measurement record screen from the QA Portal.")

    info_col, action_col = st.columns([3, 2])
    with info_col:
        st.markdown(
            """
            **Process Monitoring / Measurement Record**

            Review the existing IATF page for monthly process metrics, targets, actual values,
            and follow-up actions. This portal entry is a shortcut to the maintained Rails screen.
            """
        )
        st.markdown(
            """
            Supported use:
            - Open the current monitoring record page
            - Review annual/monthly monitoring results
            - Jump back to the main IATF product portal
            """
        )
    with action_col:
        st.link_button("Open Monitoring Record", PROCESS_MONITORING_URL, use_container_width=True)
        st.link_button("Open IATF Products", f"{IATF_BASE_URL}/products", use_container_width=True)
        st.caption("If the page does not open, check the Rails app on the configured IATF base URL.")



# --- SIDEBAR ---
st.sidebar.title("QA Toolkit")

# File Upload Section in Sidebar
st.sidebar.markdown("---")
st.sidebar.subheader("Upload Knowledge")
uploaded_file = st.sidebar.file_uploader("Add to Knowledge Base", type=["pdf", "xlsx", "docx", "pptx", "dxf", "txt"])
if uploaded_file is not None:
    if st.sidebar.button("Upload & Ingest"):
        with st.sidebar.status("Uploading..."):
            success, path, status = save_uploaded_file(uploaded_file, INGEST_DIR)
            if success:
                st.write(f"笨・{status}: {INGEST_DIR}")
                if status == "Saved": st.write("竢ｳ Ingestion started")
            else:
                st.error(f"Failed: {path}")

st.sidebar.markdown("---")

page = st.sidebar.radio("Select Tool", [
    "Home", 
    *([] if not PLATING_FLAG else ["Plating Quality Analysis"]),
    "Work Instruction Generator",
    "FMEA Editor", 
    "FTA (Fault Tree)", 
    "Why-Why Analysis",
    "Work Study",
    "Process Monitoring & Measurement",
    "Email Daily Report (P016)"
])
st.sidebar.caption("Portal に独立カードがある重複ツールは、QA Dashboard から順次整理しています。")
st.sidebar.markdown("[Open 3D Converter from Portal](http://localhost:8088/portal.html)")
st.sidebar.markdown("[Open Tolerance Center from Portal](http://localhost:8088/portal.html)")
st.sidebar.markdown("[Open Kindle Author from Portal](http://localhost:8088/portal.html)")

# --- PAGES ---

if page == "Home":
    st.title("Clawstack QA Portal")
    st.markdown(f"""
    **New Feature:**
    *   **統 Work Instruction Generator:** Upload Documents/Video/Audio to `/consume/WIP`. AI generates standard work steps.
    
    **Knowledge Base:**
    *   **Ingest:** Upload to `/consume/PFMEA_5WHY_FTA_etc`.
    *   **RAG:** Documents are indexed for FMEA/FTA analysis.
    """)
    if PLATING_FLAG:
        st.markdown("---")
        st.subheader("Plating Quality Analysis")
        card_col1, card_col2 = st.columns([3, 2])
        with card_col1:
            st.markdown(
                """
                **Plating / Reflow Analysis**

                Enter plating line conditions, reflow conditions, thermal profile, and inspection assets.
                The Portal stores case inputs, runs analysis, and keeps generated artifacts together.
                """
            )
        with card_col2:
            st.info("Use the left sidebar to open `Plating Quality Analysis` and run the analysis.")

    st.markdown("---")
    st.subheader("Process Monitoring & Measurement")
    pm_col1, pm_col2 = st.columns([3, 2])
    with pm_col1:
        st.markdown(
            """
            **IATF Process Monitoring Record**

            Open the existing process monitoring and measurement record page for KPI review,
            target-vs-actual tracking, and monthly follow-up status.
            """
        )
    with pm_col2:
        st.link_button("Open Record", PROCESS_MONITORING_URL, use_container_width=True)

elif page == "Plating Quality Analysis":
    render_plating_quality_page()

elif page == "Process Monitoring & Measurement":
    render_process_monitoring_measurement_page()

elif page == "Work Instruction Generator":
    st.header("統 Work Instruction Generator")
    st.info("Upload raw materials (PDF, Excel, Video, Audio). AI will draft a Standard Operating Procedure (SOP).")
    
    uploaded_wip = st.file_uploader("Upload Raw Material", accept_multiple_files=True)
    
    if uploaded_wip:
        if st.button("噫 Generate Instruction"):
            combined_text = ""
            media_files = []
            
            progress = st.progress(0)
            status = st.empty()
            
            for i, uf in enumerate(uploaded_wip):
                status.write(f"Processing {uf.name}...")
                success, path, stat = save_uploaded_file(uf, WIP_DIR)
                
                if success:
                    ext = os.path.splitext(path)[1].lower()
                    if ext in [".pdf", ".docx", ".xlsx", ".txt"]:
                        txt = extract_text_immediate(path)
                        combined_text += f"\n\n--- Source: {uf.name} ---\n{txt}"
                    elif ext in [".mp4", ".avi", ".mov", ".mp3", ".wav"]:
                        media_files.append(uf.name)
                        combined_text += f"\n\n--- Source: {uf.name} ---\n[Media File Present: Analysis requires Vision/Audio Module update. Using context from documents if available.]"
                
                progress.progress((i + 1) / len(uploaded_wip))
            
            status.write("ｧ AI Generating SOP...")
            
            prompt = f"""
            Create a detailed 'Work Instruction' (Standard Operating Procedure) based on the following raw materials.
            Structure it with:
            1. Title
            2. Safety Warnings (Important!)
            3. Tools Required
            4. Step-by-Step Instructions (Concrete, Action-Oriented)
            
            RAW MATERIAL CONTENT:
            {combined_text[:6000]} 
            """
            # Truncate to avoid context limit overflow if huge
            
            sop = ask_ai(prompt)
            
            col1, col2 = st.columns([2,1])
            with col1:
                st.subheader("Draft Instruction")
                st.markdown(sop)
            with col2:
                if media_files:
                    st.warning(f"Media files referenced ({len(media_files)}). AI relied on text documents for detail.")
                
                fn = f"SOP_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.md"
                if st.button("沈 Save SOP"):
                    with open(os.path.join(WORK_DIR, fn), "w", encoding="utf-8") as f:
                        f.write(sop)
                    st.success(f"Saved to {WORK_DIR}/{fn}")

elif page == "FMEA Editor":
    st.header("投 PFMEA Workspace")
    st.caption("実務向けの PFMEA 列に寄せています。ユーザーが表を作り、AI は会議メンバーとして抜け漏れと妥当性をレビューします。")
    process_step = st.text_input("Process Step", "Battery Weld")
    process_function = st.text_input("Process Function", "Join tab and terminal with stable nugget size")
    requirement = st.text_input("Requirement / Customer Need", "No crack, no leakage, electrical resistance within spec")
    st.session_state.fmea_data = normalize_pfmea_dataframe(st.session_state.get("fmea_data"), process_step)
    if st.button("Add starter PFMEA row"):
        starter_df = st.session_state.fmea_data.copy()
        starter_df.loc[len(starter_df)] = default_pfmea_rows(process_step)[0]
        st.session_state.fmea_data = normalize_pfmea_dataframe(starter_df, process_step)

    edited_df = st.data_editor(
        st.session_state.fmea_data,
        num_rows="dynamic",
        use_container_width=True,
        hide_index=True,
        column_config={
            "Process Step": st.column_config.TextColumn(width="medium"),
            "Process Function": st.column_config.TextColumn(width="medium"),
            "Requirement": st.column_config.TextColumn(width="medium"),
            "Potential Failure Mode": st.column_config.TextColumn(width="large"),
            "Potential Effect": st.column_config.TextColumn(width="large"),
            "Severity": st.column_config.NumberColumn(min_value=0, max_value=10, step=1),
            "Potential Cause": st.column_config.TextColumn(width="large"),
            "Occurrence": st.column_config.NumberColumn(min_value=0, max_value=10, step=1),
            "Current Prevention Control": st.column_config.TextColumn(width="large"),
            "Current Detection Control": st.column_config.TextColumn(width="large"),
            "Detection": st.column_config.NumberColumn(min_value=0, max_value=10, step=1),
            "RPN": st.column_config.NumberColumn(disabled=True),
            "Recommended Action": st.column_config.TextColumn(width="large"),
            "Responsibility": st.column_config.TextColumn(width="small"),
            "Due Date": st.column_config.TextColumn(width="small", help="例: 2026-04-15"),
            "Action Status": st.column_config.SelectboxColumn(options=["Open", "In Progress", "Done", "Hold"]),
        },
    )
    edited_df = normalize_pfmea_dataframe(edited_df, process_step)
    edited_df["Process Function"] = edited_df["Process Function"].replace("", process_function)
    edited_df["Requirement"] = edited_df["Requirement"].replace("", requirement)
    st.session_state.fmea_data = edited_df

    risk_col1, risk_col2, risk_col3, risk_col4 = st.columns(4)
    risk_col1.metric("Rows", len(edited_df))
    risk_col2.metric("High RPN Rows", int((edited_df["RPN"] >= 100).sum()))
    risk_col3.metric("Open Actions", int((edited_df["Action Status"] != "Done").sum()))
    risk_col4.metric("Max RPN", int(edited_df["RPN"].max()) if len(edited_df) else 0)

    with st.expander("PFMEA Review Focus"):
        st.markdown(
            "- `Potential Failure Mode` が現象になっているか\n"
            "- `Potential Effect` が顧客・次工程影響まで届いているか\n"
            "- `Potential Cause` が真因候補になっているか\n"
            "- `Current Prevention / Detection Control` が分かれているか\n"
            "- `Recommended Action` と `Responsibility / Due Date` が会議で決められる粒度か"
        )

    mode_name = st.selectbox(
        "FMEA review mode",
        list(WHYWHY_MODE_PRESETS.keys()),
        index=1,
        key="fmea_mode_name",
    )
    mode_preset = WHYWHY_MODE_PRESETS[mode_name]
    st.caption(mode_preset["description"])

    setup_col1, setup_col2 = st.columns([2, 1])
    with setup_col1:
        selected_agents = st.multiselect(
            "FMEA AI participants",
            options=list(WHYWHY_AGENT_CATALOG.keys()),
            default=["Gemini", "ChatGPT"],
            max_selections=3,
            key="fmea_agents",
        )
    with setup_col2:
        st.markdown("**Current role plan**")
        for idx, agent_name in enumerate(selected_agents or ["Gemini", "ChatGPT"], start=1):
            role = ["Lead Investigator", "Logic Auditor", "Countermeasure Critic"][idx - 1]
            st.caption(f"{agent_name}: {role}")

    options_col1, options_col2 = st.columns(2)
    with options_col1:
        use_internal_docs = st.checkbox("Use Mitsui / internal RAG", value=True, key="fmea_internal_docs")
    with options_col2:
        use_web_docs = st.checkbox("Use Web knowledge", value=mode_preset["use_web"], key="fmea_web_docs")

    effective_agents = (selected_agents or ["Gemini", "ChatGPT"])[: mode_preset["max_agents"]]
    effective_use_web = use_web_docs and mode_preset["use_web"]
    st.info(
        " | ".join([
            f"Participants: {', '.join(effective_agents)}",
            f"Web knowledge: {'On' if effective_use_web else 'Off'}",
            f"Consensus synthesis: {'On' if mode_preset['synthesis'] else 'Off'}",
            f"Internal RAG top-k: {mode_preset['top_k']}",
        ])
    )
    st.caption(
        f"Deep model setting: `{FMEA_DEEP_MODEL}` via `{LITELLM_URL}/chat/completions`. "
        "単独の軽い確認が必要な場合だけ Local Ollama へ戻すより、まずは省APIモードを使うのがおすすめです。"
    )

    if st.button("Run Multi-AI FMEA Review"):
        effective_agents = effective_agents or ["Gemini"]
        records = []
        for row in edited_df.fillna("").to_dict(orient="records"):
            step_value = row.get("Process Step") or process_step
            records.append(
                f"- Process Step: {step_value} | Function: {row.get('Process Function', '')} | "
                f"Requirement: {row.get('Requirement', '')} | Failure Mode: {row.get('Potential Failure Mode', '')} | "
                f"Effect: {row.get('Potential Effect', '')} | Cause: {row.get('Potential Cause', '')} | "
                f"S/O/D: {row.get('Severity', '')}/{row.get('Occurrence', '')}/{row.get('Detection', '')} | "
                f"RPN: {row.get('RPN', '')} | Prevention: {row.get('Current Prevention Control', '')} | "
                f"Detection Control: {row.get('Current Detection Control', '')} | Action: {row.get('Recommended Action', '')} | "
                f"Owner/Due: {row.get('Responsibility', '')}/{row.get('Due Date', '')} | Status: {row.get('Action Status', '')}"
            )
        fmea_body = (
            f"Process Step: {process_step}\n"
            f"Process Function: {process_function}\n"
            f"Requirement: {requirement}\n"
            f"Current PFMEA rows:\n" + ("\n".join(records) if records else f"- Process Step: {process_step}")
        )
        with st.spinner("Reviewing FMEA with multiple AI agents..."):
            query, rag_context, web_context, context = build_standard_review_context(
                query=f"PFMEA {process_step} {process_function} {requirement} {' '.join([str(v) for v in edited_df.fillna('').astype(str).values.flatten()[:20]])}",
                web_query=f"PFMEA manufacturing process function failure mode effect cause controls action {process_step}",
                use_internal_docs=use_internal_docs,
                use_web_docs=effective_use_web,
                rag_limit=mode_preset["top_k"],
            )
            agent_outputs, synthesis, usage_rollup = run_multi_agent_quality_review(
                subject_title="FMEA draft",
                subject_body=fmea_body,
                selected_agent_names=effective_agents,
                context_text=context,
                review_task_text=(
                    "1. Review the PFMEA rows for logical quality and real manufacturing usability.\n"
                    "2. Check whether Process Function, Requirement, Failure Mode, Effect, and Cause are properly separated.\n"
                    "3. Comment on whether S/O/D severity appears justified.\n"
                    "4. Point out missing prevention controls, missing detection controls, and weak recommended actions.\n"
                    "5. Suggest where FTA or Why-Why deepening is needed for high-risk causes.\n"
                    "6. Highlight what internal standards or procedures should be checked next.\n"
                    "7. Keep the response concise and structured for a PFMEA meeting."
                ),
                synthesis_task_text=(
                    "Create a final integrated PFMEA review with these sections:\n"
                    "1. Consensus Summary\n"
                    "2. Recommended PFMEA row corrections\n"
                    "3. S/O/D concerns\n"
                    "4. Missing controls or actions\n"
                    "5. Where FTA or Why-Why should be used next\n"
                    "6. Evidence / procedures to confirm next"
                ),
                synthesis_enabled=mode_preset["synthesis"],
            )

        st.subheader("Integrated Conclusion")
        st.markdown(synthesis)

        usage_col1, usage_col2 = st.columns([2, 1])
        with usage_col1:
            st.subheader("Turn Usage")
            st.caption(format_usage_summary(usage_rollup))
        with usage_col2:
            st.metric("Deep AI Calls", usage_rollup.get("calls", 0))
            st.metric("Total Tokens", f"{usage_rollup.get('total_tokens', 0):,}")

        st.subheader("Agent Opinions")
        for item in agent_outputs:
            with st.expander(f"{item['agent_name']} | {item['assigned_role']} | {item['model']}"):
                st.caption(item["focus"])
                st.markdown(item["reply"])
                st.caption(format_usage_summary({
                    "prompt_tokens": item.get("usage", {}).get("prompt_tokens", 0),
                    "completion_tokens": item.get("usage", {}).get("completion_tokens", 0),
                    "total_tokens": item.get("usage", {}).get("total_tokens", 0),
                    "estimated_cost_usd": item.get("estimated_cost_usd"),
                }))

        if rag_context or web_context:
            st.subheader("References")
            if rag_context:
                with st.expander("Internal / PDF Knowledge"):
                    st.markdown(rag_context[:2000])
            if web_context:
                with st.expander("Web Knowledge"):
                    st.markdown(web_context[:2000])
            st.caption(f"Search query: {query}")

elif page == "FTA (Fault Tree)":
    st.header("元 Fault Tree")
    st.caption("FTA も Why-Why と同じく、既定 2 名の AI で原因展開をチェックできるようにしています。")
    top_event = st.text_input("Top Event", "Motor Stall")
    nodes = st.text_area("Define Causes (Lines)", "Overload\nShort Circuit").split('\n')
    mermaid = f"graph TD\nTOP[\"{top_event}\"] --> OR((OR))"
    for i, n in enumerate(nodes):
        if n.strip(): mermaid += f"\nOR --> C{i}[\"{n.strip()}\"]"
    st.mermaid(mermaid)

    mode_name = st.selectbox(
        "FTA review mode",
        list(WHYWHY_MODE_PRESETS.keys()),
        index=1,
        key="fta_mode_name",
    )
    mode_preset = WHYWHY_MODE_PRESETS[mode_name]
    st.caption(mode_preset["description"])

    setup_col1, setup_col2 = st.columns([2, 1])
    with setup_col1:
        selected_agents = st.multiselect(
            "FTA AI participants",
            options=list(WHYWHY_AGENT_CATALOG.keys()),
            default=["Gemini", "ChatGPT"],
            max_selections=3,
            key="fta_agents",
        )
    with setup_col2:
        st.markdown("**Current role plan**")
        for idx, agent_name in enumerate(selected_agents or ["Gemini", "ChatGPT"], start=1):
            role = ["Lead Investigator", "Logic Auditor", "Countermeasure Critic"][idx - 1]
            st.caption(f"{agent_name}: {role}")

    options_col1, options_col2 = st.columns(2)
    with options_col1:
        use_internal_docs = st.checkbox("Use Mitsui / internal RAG", value=True, key="fta_internal_docs")
    with options_col2:
        use_web_docs = st.checkbox("Use Web knowledge", value=mode_preset["use_web"], key="fta_web_docs")

    effective_agents = (selected_agents or ["Gemini", "ChatGPT"])[: mode_preset["max_agents"]]
    effective_use_web = use_web_docs and mode_preset["use_web"]
    st.info(
        " | ".join([
            f"Participants: {', '.join(effective_agents)}",
            f"Web knowledge: {'On' if effective_use_web else 'Off'}",
            f"Consensus synthesis: {'On' if mode_preset['synthesis'] else 'Off'}",
            f"Internal RAG top-k: {mode_preset['top_k']}",
        ])
    )

    if st.button("Run Multi-AI FTA Review"):
        effective_agents = effective_agents or ["Gemini"]
        cause_lines = [item.strip() for item in nodes if item.strip()]
        fta_body = (
            f"Top event: {top_event}\n"
            f"Current fault tree branch draft:\n" +
            "\n".join([f"- {item}" for item in cause_lines])
        )
        with st.spinner("Reviewing FTA with multiple AI agents..."):
            query, rag_context, web_context, context = build_standard_review_context(
                query=f"FTA {top_event} {' '.join(cause_lines)}",
                web_query=f"fault tree analysis manufacturing root causes {top_event}",
                use_internal_docs=use_internal_docs,
                use_web_docs=effective_use_web,
                rag_limit=mode_preset["top_k"],
            )
            agent_outputs, synthesis, usage_rollup = run_multi_agent_quality_review(
                subject_title="FTA draft",
                subject_body=fta_body,
                selected_agent_names=effective_agents,
                context_text=context,
                review_task_text=(
                    "1. Review whether the listed causes are logically connected to the top event.\n"
                    "2. Suggest missing intermediate causes, branch separation, or gate logic concerns.\n"
                    "3. Point out if immediate causes and root causes are mixed together.\n"
                    "4. Identify what evidence, records, or procedures should be checked next.\n"
                    "5. Keep the response concise and structured for an FTA review meeting."
                ),
                synthesis_task_text=(
                    "Create a final integrated FTA review with these sections:\n"
                    "1. Consensus Summary\n"
                    "2. Recommended branch structure\n"
                    "3. Missing intermediate or root causes\n"
                    "4. Gate logic / tree quality concerns\n"
                    "5. Evidence to confirm next"
                ),
                synthesis_enabled=mode_preset["synthesis"],
            )

        st.subheader("Integrated Conclusion")
        st.markdown(synthesis)

        usage_col1, usage_col2 = st.columns([2, 1])
        with usage_col1:
            st.subheader("Turn Usage")
            st.caption(format_usage_summary(usage_rollup))
        with usage_col2:
            st.metric("Deep AI Calls", usage_rollup.get("calls", 0))
            st.metric("Total Tokens", f"{usage_rollup.get('total_tokens', 0):,}")

        st.subheader("Agent Opinions")
        for item in agent_outputs:
            with st.expander(f"{item['agent_name']} | {item['assigned_role']} | {item['model']}"):
                st.caption(item["focus"])
                st.markdown(item["reply"])
                st.caption(format_usage_summary({
                    "prompt_tokens": item.get("usage", {}).get("prompt_tokens", 0),
                    "completion_tokens": item.get("usage", {}).get("completion_tokens", 0),
                    "total_tokens": item.get("usage", {}).get("total_tokens", 0),
                    "estimated_cost_usd": item.get("estimated_cost_usd"),
                }))

        if rag_context or web_context:
            st.subheader("References")
            if rag_context:
                with st.expander("Internal / PDF Knowledge"):
                    st.markdown(rag_context[:2000])
            if web_context:
                with st.expander("Web Knowledge"):
                    st.markdown(web_context[:2000])
            st.caption(f"Search query: {query}")

elif page == "Why-Why Analysis":
    st.header("5-Whys (Multi-AI Review)")
    st.caption("既定は 2 名の AI で討議し、必要なら 3 名まで増やせます。社内文書 RAG を優先し、Web 知識は補助参照にします。")
    problem = st.text_input("Problem", "Leakage")
    whys = [st.text_input(f"{i}. Why?", key=f"w{i}") for i in range(1, 6)]
    mode_name = st.selectbox(
        "Review mode",
        list(WHYWHY_MODE_PRESETS.keys()),
        index=1,
    )
    mode_preset = WHYWHY_MODE_PRESETS[mode_name]
    st.caption(mode_preset["description"])

    setup_col1, setup_col2 = st.columns([2, 1])
    with setup_col1:
        selected_agents = st.multiselect(
            "AI participants",
            options=list(WHYWHY_AGENT_CATALOG.keys()),
            default=["Gemini", "ChatGPT"],
            max_selections=3,
            help="既定は 2 名です。Gemini / ChatGPT / Claude から選べます。",
        )
    with setup_col2:
        st.markdown("**Current role plan**")
        for idx, agent_name in enumerate(selected_agents or ["Gemini", "ChatGPT"], start=1):
            role = ["Lead Investigator", "Logic Auditor", "Countermeasure Critic"][idx - 1]
            st.caption(f"{agent_name}: {role}")

    options_col1, options_col2 = st.columns(2)
    with options_col1:
        use_internal_docs = st.checkbox("Use Mitsui / internal RAG", value=True)
    with options_col2:
        use_web_docs = st.checkbox("Use Web knowledge", value=mode_preset["use_web"])

    effective_agents = (selected_agents or ["Gemini", "ChatGPT"])[: mode_preset["max_agents"]]
    effective_use_web = use_web_docs and mode_preset["use_web"]
    runtime_notes = [
        f"Participants: {', '.join(effective_agents)}",
        f"Web knowledge: {'On' if effective_use_web else 'Off'}",
        f"Consensus synthesis: {'On' if mode_preset['synthesis'] else 'Off'}",
        f"Internal RAG top-k: {mode_preset['top_k']}",
        "Backward validation: On",
    ]
    st.info(" | ".join(runtime_notes))

    if st.button("Run Multi-AI Review"):
        effective_agents = effective_agents or ["Gemini"]
        with st.spinner("Reviewing with multiple AI agents..."):
            query, rag_context, web_context, context = build_whywhy_context(
                problem,
                whys,
                use_internal_docs=use_internal_docs,
                use_web_docs=effective_use_web,
                rag_limit=mode_preset["top_k"],
            )
            agent_outputs, synthesis, usage_rollup = run_whywhy_agents(
                problem,
                whys,
                effective_agents,
                context,
                synthesis_enabled=mode_preset["synthesis"],
            )

        st.subheader("Integrated Conclusion")
        st.markdown(synthesis)

        usage_col1, usage_col2 = st.columns([2, 1])
        with usage_col1:
            st.subheader("Turn Usage")
            st.caption(format_usage_summary(usage_rollup))
        with usage_col2:
            st.metric("Deep AI Calls", usage_rollup.get("calls", 0))
            st.metric("Total Tokens", f"{usage_rollup.get('total_tokens', 0):,}")

        st.subheader("Agent Opinions")
        for item in agent_outputs:
            with st.expander(f"{item['agent_name']} | {item['assigned_role']} | {item['model']}"):
                st.caption(item["focus"])
                st.markdown(item["reply"])
                st.caption(format_usage_summary({
                    "prompt_tokens": item.get("usage", {}).get("prompt_tokens", 0),
                    "completion_tokens": item.get("usage", {}).get("completion_tokens", 0),
                    "total_tokens": item.get("usage", {}).get("total_tokens", 0),
                    "estimated_cost_usd": item.get("estimated_cost_usd"),
                }))

        if rag_context or web_context:
            st.subheader("References")
            if rag_context:
                with st.expander("Internal / PDF Knowledge"):
                    st.markdown(rag_context[:2000])
            if web_context:
                with st.expander("Web Knowledge"):
                    st.markdown(web_context[:2000])
            st.caption(f"Search query: {query}")

elif page == "Work Study":
    st.header("竢ｱ・・Work Study")
    uploaded_vid = st.file_uploader("Upload Video", type=["mp4", "avi"])
    if uploaded_vid:
        success, path, _ = save_uploaded_file(uploaded_vid, WORK_DIR)
        if success: st.success(f"Video ready at {path}")

elif page == "3D Converter":
    st.header("3D Converter")
    st.markdown("""
    **2D / 3D conversion tools**
    
    | Conversion | Input | Output | Usage |
    |------|------|------|------|
    | **DXF -> STEP/STL** | DXF | STEP or STL | Simple extrusion-based 3D generation |
    | **Model -> 3D HTML** | STEP/STL/OBJ | Interactive HTML | Browser preview and sharing |
    | **Model -> 3D PDF** | STEP/STL/OBJ | 3D PDF | Adobe Acrobat Reader compatible |
    
    ---
    """)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Upload File")
        conv_type = st.radio("Conversion Type", ["DXF -> STEP/STL", "Model -> 3D HTML", "Model -> 3D PDF"])
        
        if conv_type == "DXF -> STEP/STL":
            uploaded_3d = st.file_uploader("DXF file", type=["dxf"], key="dxf_upload")
            height = st.number_input("Extrusion height (mm)", min_value=0.1, value=10.0, step=0.5)
            output_format = st.selectbox("Output format", ["STEP", "STL"])
        else:
            uploaded_3d = st.file_uploader("3D model", type=["step", "stp", "stl", "obj"], key="model_upload")
            html_profile = "storage_5mb"
            html_purpose = "普通の生成"
            drawing_pdf = None
            if conv_type == "Model -> 3D HTML":
                html_purpose = st.radio(
                    "Generation Purpose",
                    ["普通の生成", "GD&T 2D/3D 整合"],
                    horizontal=True,
                    help="GD&T 2D/3D 整合を選ぶと、通常のHTMLに加えてレビュー用マニフェストとチェックリストを含む版を生成します。",
                )
                html_profile = st.selectbox(
                    "HTML Size Profile",
                    ["email_2mb", "storage_5mb", "high_quality"],
                    index=1,
                    format_func=lambda value: {
                        "email_2mb": "Email Attachment (Max 2MB)",
                        "storage_5mb": "Computer Storage (Max 5MB)",
                        "high_quality": "High Resolution (No Reduction)",
                    }[value],
                )
                if uploaded_3d is not None:
                    with tempfile.TemporaryDirectory() as td_est:
                        estimate_input_path = os.path.join(td_est, uploaded_3d.name)
                        with open(estimate_input_path, "wb") as f:
                            f.write(uploaded_3d.getbuffer())
                        estimate_cmd = [
                            "python3",
                            "/work/scripts/model2html.py",
                            estimate_input_path,
                            os.path.join(td_est, "estimate.html"),
                            "--profile",
                            html_profile,
                            "--estimate-only",
                        ]
                        estimate_result = subprocess.run(estimate_cmd, capture_output=True, text=True, timeout=300)
                        if estimate_result.returncode == 0:
                            estimate_data = json.loads(estimate_result.stdout)
                            st.caption(
                                f"Safe HTML estimate: {estimate_data.get('actual_html_kb', '--'):,} KB | "
                                f"ZIP estimate: {estimate_data.get('estimated_zip_kb', '--'):,} KB"
                            )
                            if estimate_data.get("selected_profile_ok"):
                                st.success(
                                    f"Selected profile should fit: {estimate_data.get('profile_label')} "
                                    f"({estimate_data.get('actual_html_kb', '--'):,} KB est.)"
                                )
                            else:
                                st.warning(estimate_data.get("selected_profile_error", "Selected profile may not fit."))
                        else:
                            st.warning(f"Size estimate failed: {estimate_result.stderr or estimate_result.stdout}")
                if html_purpose == "GD&T 2D/3D 整合":
                    drawing_pdf = st.file_uploader(
                        "2D Drawing PDF (Optional but recommended)",
                        type=["pdf"],
                        key="gdt_drawing_upload",
                        help="図面PDFを添付すると、レビュー bundle に抽出プレビューと候補 requirement を含めます。",
                    )
                    st.info(
                        "GD&T review mode: HTML viewer に加えて、review manifest と checklist を含むレビュー版を生成します。"
                    )
                    if drawing_pdf is not None:
                        with tempfile.TemporaryDirectory() as td_drawing_preview:
                            preview_pdf_path = os.path.join(td_drawing_preview, drawing_pdf.name)
                            with open(preview_pdf_path, "wb") as f:
                                f.write(drawing_pdf.getbuffer())
                            drawing_preview = extract_gdt_pdf_review_context(preview_pdf_path)
                        mapping_key = f"gdt_requirement_mapping::{safe_stem(uploaded_3d.name)}::{safe_stem(drawing_pdf.name)}"
                        default_mapping_rows = build_gdt_requirement_mapping_rows(
                            drawing_preview.get("candidate_requirements", [])
                        )
                        prior_rows = st.session_state.get(mapping_key)
                        if prior_rows and isinstance(prior_rows, list) and len(prior_rows) == len(default_mapping_rows):
                            default_mapping_rows = prior_rows
                        if drawing_preview.get("candidate_requirements"):
                            st.caption(
                                f"Drawing preview: {drawing_preview.get('page_count', 0)} page(s) | "
                                f"{len(drawing_preview.get('candidate_requirements', []))} likely GD&T lines detected"
                            )
                            with st.expander("Drawing extraction preview"):
                                st.write(drawing_preview.get("candidate_requirements", []))
                                if drawing_preview.get("text_preview"):
                                    st.text_area(
                                        "Preview text",
                                        drawing_preview.get("text_preview", ""),
                                        height=180,
                                        disabled=True,
                                        key="gdt_drawing_preview_text",
                                    )
                            st.markdown("**Requirement -> face/axis review table**")
                            edited_mapping_df = st.data_editor(
                                pd.DataFrame(default_mapping_rows),
                                hide_index=True,
                                use_container_width=True,
                                num_rows="fixed",
                                key=f"{mapping_key}::editor",
                                column_config={
                                    "requirement_id": st.column_config.TextColumn("Requirement ID", disabled=True, width="small"),
                                    "requirement_text": st.column_config.TextColumn("Requirement", disabled=True, width="large"),
                                    "candidate_face_ids": st.column_config.TextColumn("Candidate Face IDs"),
                                    "candidate_axis_ids": st.column_config.TextColumn("Candidate Axis IDs"),
                                    "chosen_target": st.column_config.TextColumn("Chosen Target"),
                                    "status": st.column_config.SelectboxColumn(
                                        "Status",
                                        options=["pending", "candidate_listed", "chosen", "rejected"],
                                        required=True,
                                    ),
                                    "review_note": st.column_config.TextColumn("Review Note", width="large"),
                                },
                            )
                            st.session_state[mapping_key] = edited_mapping_df.to_dict("records")
                        else:
                            st.caption(
                                f"Drawing preview: {drawing_preview.get('page_count', 0)} page(s) | "
                                "No strong GD&T keywords detected automatically"
                            )
    
    with col2:
        st.subheader("Run Conversion")
        if st.button("Start Conversion", use_container_width=True):
            if uploaded_3d is None:
                st.error("Please upload a file first.")
            else:
                import subprocess
                import tempfile
                
                with tempfile.TemporaryDirectory() as td:
                    # Save uploaded file
                    input_path = os.path.join(td, uploaded_3d.name)
                    with open(input_path, "wb") as f:
                        f.write(uploaded_3d.getbuffer())
                    
                    try:
                        if conv_type == "DXF -> STEP/STL":
                            ext = "step" if output_format == "STEP" else "stl"
                            output_path = os.path.join(td, f"output.{ext}")
                            cmd = ["python3", "/work/scripts/dxf23d.py", input_path, output_path, "--height", str(height)]
                        elif conv_type == "Model -> 3D HTML":
                            output_basename = "output_gdt_review.html" if html_purpose == "GD&T 2D/3D 整合" else "output.html"
                            output_path = os.path.join(td, output_basename)
                            cmd = [
                                "python3",
                                "/work/scripts/model2html.py",
                                input_path,
                                output_path,
                                "--profile",
                                html_profile,
                            ]
                        else:  # 3D PDF
                            output_path = os.path.join(td, "output.pdf")
                            cmd = ["python3", "/work/scripts/model2pdf.py", input_path, output_path]
                        
                        with st.spinner("Converting..."):
                            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
                        
                        if result.returncode != 0:
                            st.error(f"Conversion failed: {result.stderr or result.stdout}")
                        elif os.path.exists(output_path):
                            st.success("Conversion completed.")
                            if conv_type == "Model -> 3D HTML":
                                gdt_manifest = None
                                drawing_info = {}
                                drawing_pdf_name = ""
                                drawing_pdf_bytes = None
                                if html_purpose == "GD&T 2D/3D 整合" and drawing_pdf is not None:
                                    drawing_pdf_name = drawing_pdf.name
                                    drawing_pdf_bytes = drawing_pdf.getbuffer().tobytes()
                                    drawing_pdf_path = os.path.join(td, drawing_pdf_name)
                                    with open(drawing_pdf_path, "wb") as f:
                                        f.write(drawing_pdf_bytes)
                                    drawing_info = extract_gdt_pdf_review_context(drawing_pdf_path)
                                requirement_mappings = build_gdt_requirement_mapping_rows(
                                    drawing_info.get("candidate_requirements", [])
                                )
                                if html_purpose == "GD&T 2D/3D 整合" and drawing_pdf is not None:
                                    mapping_key = f"gdt_requirement_mapping::{safe_stem(uploaded_3d.name)}::{safe_stem(drawing_pdf.name)}"
                                    stored_mapping = st.session_state.get(mapping_key)
                                    if stored_mapping and isinstance(stored_mapping, list):
                                        requirement_mappings = stored_mapping
                                if html_purpose == "GD&T 2D/3D 整合":
                                    gdt_manifest = build_gdt_review_manifest(
                                        uploaded_3d.name,
                                        html_profile,
                                        drawing_info=drawing_info,
                                        requirement_mappings=requirement_mappings,
                                    )
                                with open(output_path, "rb") as f:
                                    html_bytes = f.read()
                                output_name = os.path.basename(output_path)
                                zip_name = os.path.splitext(output_name)[0] + ".zip"
                                zip_buffer = io.BytesIO()
                                with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
                                    zf.writestr(output_name, html_bytes)
                                    if html_purpose == "GD&T 2D/3D 整合":
                                        checklist_md = build_gdt_checklist_markdown(uploaded_3d.name, drawing_name=drawing_info.get("drawing_name", ""))
                                        prefix = safe_stem(uploaded_3d.name)
                                        zf.writestr(
                                            f"{prefix}_gdt_review_manifest.json",
                                            json.dumps(gdt_manifest, ensure_ascii=False, indent=2),
                                        )
                                        zf.writestr(f"{prefix}_gdt_checklist.md", checklist_md)
                                        if drawing_pdf_name and drawing_pdf_bytes is not None:
                                            zf.writestr(f"{prefix}_source_drawing.pdf", drawing_pdf_bytes)
                                zip_bytes = zip_buffer.getvalue()
                                dl_html, dl_zip = st.columns(2)
                                with dl_html:
                                    st.download_button(
                                        "Download 3D HTML" if html_purpose == "普通の生成" else "Download GD&T HTML",
                                        html_bytes,
                                        file_name=output_name,
                                        mime="text/html",
                                        use_container_width=True,
                                    )
                                with dl_zip:
                                    st.download_button(
                                        "Download ZIP" if html_purpose == "普通の生成" else "Download GD&T Review ZIP",
                                        zip_bytes,
                                        file_name=zip_name,
                                        mime="application/zip",
                                        use_container_width=True,
                                    )
                                st.caption(
                                    f"The downloaded HTML can be opened locally by double-clicking it in a browser. "
                                    f"HTML size: {round(len(html_bytes) / 1024):,} KB | ZIP size: {round(len(zip_bytes) / 1024):,} KB"
                                )
                                if html_purpose == "GD&T 2D/3D 整合":
                                    st.markdown("**GD&T review bundle contents**")
                                    st.json(gdt_manifest)
                                    if drawing_info.get("candidate_requirements"):
                                        mapping_summary = summarize_gdt_requirement_mappings(requirement_mappings)
                                        mcol1, mcol2, mcol3, mcol4 = st.columns(4)
                                        with mcol1:
                                            st.metric("Requirements", mapping_summary["total_requirements"])
                                        with mcol2:
                                            st.metric("Chosen", mapping_summary["chosen_count"])
                                        with mcol3:
                                            st.metric("Candidate Listed", mapping_summary["candidate_listed_count"])
                                        with mcol4:
                                            st.metric("Pending", mapping_summary["pending_count"])
                                        st.markdown("**Drawing-derived candidate requirements**")
                                        st.write(drawing_info.get("candidate_requirements", []))
                                        st.markdown("**Saved requirement mapping rows**")
                                        st.dataframe(pd.DataFrame(requirement_mappings), use_container_width=True, hide_index=True)
                                    st.caption(
                                        "This bundle is review-oriented. Exact face ids / chosen targets still need engineer confirmation against the 2D drawing."
                                    )
                            else:
                                preview_path = os.path.join(
                                    td,
                                    os.path.splitext(os.path.basename(output_path))[0] + "_outline_preview.pdf",
                                )
                                dl1, dl2 = st.columns(2)
                                with dl1:
                                    with open(output_path, "rb") as f:
                                        st.download_button(
                                            "Download 3D PDF",
                                            f.read(),
                                            file_name=os.path.basename(output_path),
                                            use_container_width=True
                                        )
                                with dl2:
                                    if os.path.exists(preview_path):
                                        with open(preview_path, "rb") as f:
                                            st.download_button(
                                                "Download Outline Preview",
                                                f.read(),
                                                file_name=os.path.basename(preview_path),
                                                use_container_width=True
                                            )
                        else:
                            st.error("Output file was not generated.")
                    except subprocess.TimeoutExpired:
                        st.error("Conversion timed out after 5 minutes.")
                    except Exception as e:
                        st.error(f"Error: {e}")
    
    st.markdown("---")
    st.info("""
    **Notes**
    - **DXF -> STEP/STL:** best for simple contour-based extrusion workflows.
    - **3D HTML:** uses browser-friendly interactive preview output.
    - **3D PDF:** creates Acrobat-compatible 3D PDF from STEP/STL/OBJ.
    """)

# -------------------------
# 蜈ｬ蟾ｮ隗｣譫舌・繝ｼ繧ｸ (Cetol6Sigma Style with 3D Viewer)
# -------------------------
elif page == "Tolerance Analysis":
    st.header("Tolerance Analysis Tool (Cetol6Sigma Style)")
    
    # Initialize session state
    if 'tol_dimensions' not in st.session_state:
        st.session_state.tol_dimensions = []
    if 'tol_result' not in st.session_state:
        st.session_state.tol_result = None
    if 'mesh_data' not in st.session_state:
        st.session_state.mesh_data = None
    if 'extracted_dims' not in st.session_state:
        st.session_state.extracted_dims = []
    
    # Tab layout
    tab_model, tab_dims, tab_result = st.tabs(["塙 3D繝｢繝・Ν & 謚ｽ蜃ｺ", "笵難ｸ・蜈ｬ蟾ｮ繝√ぉ繝ｼ繝ｳ", "投 隗｣譫千ｵ先棡"])
    
    # === Tab 1: 3D Model & Extraction ===
    with tab_model:
        col_upload, col_3d = st.columns([1, 2])
        
        with col_upload:
            st.subheader("Upload CAD File")
            uploaded_cad = st.file_uploader("STEP/STL 繝輔ぃ繧､繝ｫ", type=["step", "stp", "stl"], key="tol_upload")
            
            if uploaded_cad:
                st.success(f"笨・{uploaded_cad.name} ({len(uploaded_cad.getvalue())/1024:.1f} KB)")
                default_tolerance = st.number_input("繝・ヵ繧ｩ繝ｫ繝亥・蟾ｮ (ﾂｱmm)", min_value=0.001, value=0.1, step=0.01)
                
                if st.button("剥 3D繝｢繝・Ν隱ｭ縺ｿ霎ｼ縺ｿ & 蟇ｸ豕墓歓蜃ｺ", use_container_width=True, type="primary"):
                    import tempfile
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_cad.name)[1]) as tf:
                        tf.write(uploaded_cad.getbuffer())
                        temp_path = tf.name
                    
                    with st.spinner("繝｡繝・す繝･謚ｽ蜃ｺ荳ｭ..."):
                        try:
                            ext = os.path.splitext(uploaded_cad.name)[1].lower()
                            
                            if ext == ".stl":
                                # Use trimesh for STL (available in quality_dashboard)
                                import trimesh
                                mesh = trimesh.load(temp_path)
                                
                                st.session_state.mesh_data = {
                                    "vertices": mesh.vertices.tolist(),
                                    "faces": mesh.faces.tolist(),
                                    "face_colors": ["#667eea"] * len(mesh.faces)
                                }
                                
                                # Extract bounding box dimensions
                                size = mesh.bounds[1] - mesh.bounds[0]
                                st.session_state.extracted_dims = [
                                    {"id": "X", "name": f"BBox X", "nominal": round(float(size[0]), 4), "tolerance": default_tolerance, "label": "X"},
                                    {"id": "Y", "name": f"BBox Y", "nominal": round(float(size[1]), 4), "tolerance": default_tolerance, "label": "Y"},
                                    {"id": "Z", "name": f"BBox Z", "nominal": round(float(size[2]), 4), "tolerance": default_tolerance, "label": "Z"},
                                ]
                                st.success("STL mesh extraction completed.")
                                
                            else:
                                # For STEP, use docker exec to antigravity with FreeCAD
                                docker_temp = f"/tmp/cad_input{ext}"
                                
                                # Copy file to container
                                copy_cmd = f'docker cp "{temp_path}" clawstack-antigravity-1:{docker_temp}'
                                result = subprocess.run(copy_cmd, shell=True, capture_output=True, text=True)
                                
                                if result.returncode != 0:
                                    st.error(f"繝輔ぃ繧､繝ｫ繧ｳ繝斐・繧ｨ繝ｩ繝ｼ: {result.stderr}")
                                else:
                                    # Run extraction script
                                    extract_cmd = f'docker exec clawstack-antigravity-1 python3 /work/scripts/extract_mesh.py {docker_temp}'
                                    result = subprocess.run(extract_cmd, shell=True, capture_output=True, text=True, timeout=120)
                                    
                                    if result.returncode == 0 and result.stdout:
                                        mesh_result = json.loads(result.stdout)
                                        if "error" in mesh_result:
                                            st.error(f"謚ｽ蜃ｺ繧ｨ繝ｩ繝ｼ: {mesh_result['error']}")
                                        else:
                                            st.session_state.mesh_data = mesh_result
                                            st.session_state.extracted_dims = mesh_result.get("dimensions", [])
                                            for d in st.session_state.extracted_dims:
                                                d["tolerance"] = default_tolerance
                                            st.success(f"笨・{len(st.session_state.extracted_dims)} 蟇ｸ豕輔ｒ謚ｽ蜃ｺ")
                                    else:
                                        st.warning("STEP extraction failed. Falling back to external mesh mode.")
                                        st.session_state.mesh_data = None
                            
                        except Exception as e:
                            st.error(f"蜃ｦ逅・お繝ｩ繝ｼ: {e}")
                        finally:
                            os.unlink(temp_path)
                    
                    st.rerun()
            
            # Extracted dimensions list
            if st.session_state.extracted_dims:
                st.markdown("---")
                st.subheader("Extracted Dimensions")
                for i, dim in enumerate(st.session_state.extracted_dims[:10]):  # Show first 10
                    col_d1, col_d2 = st.columns([3, 1])
                    with col_d1:
                        label = dim.get("label", chr(65 + i))
                        st.write(f"**{label}**: {dim['name']} = {dim['nominal']:.3f} mm")
                    with col_d2:
                        if st.button("Add", key=f"add_dim_{i}", help="Add this dimension to the chain"):
                            st.session_state.tol_dimensions.append({
                                "name": dim['name'],
                                "nominal": dim['nominal'],
                                "tolerance": dim.get('tolerance', 0.1),
                                "direction": "+"
                            })
                            st.rerun()
        
        with col_3d:
            st.subheader("塙 3D繝薙Η繝ｼ")
            
            if st.session_state.mesh_data and "vertices" in st.session_state.mesh_data:
                import plotly.graph_objects as go
                import numpy as np
                
                mesh = st.session_state.mesh_data
                verts = np.array(mesh["vertices"])
                faces = np.array(mesh["faces"])
                
                # Create mesh3d trace
                fig = go.Figure(data=[
                    go.Mesh3d(
                        x=verts[:, 0],
                        y=verts[:, 1],
                        z=verts[:, 2],
                        i=faces[:, 0],
                        j=faces[:, 1],
                        k=faces[:, 2],
                        color='#667eea',
                        opacity=0.8,
                        flatshading=True,
                        lighting=dict(ambient=0.5, diffuse=0.8, specular=0.3),
                        lightposition=dict(x=100, y=100, z=100),
                        hoverinfo='text',
                        hovertext='Click mesh faces to inspect dimensions'
                    )
                ])
                
                # Add face labels if available
                if "bounding_box" in mesh and mesh["bounding_box"]:
                    bbox = mesh["bounding_box"]
                    center = [(bbox["min"][i] + bbox["max"][i]) / 2 for i in range(3)]
                    
                    # Add dimension annotations
                    fig.add_trace(go.Scatter3d(
                        x=[bbox["min"][0], bbox["max"][0]],
                        y=[center[1], center[1]],
                        z=[bbox["min"][2], bbox["min"][2]],
                        mode='lines+text',
                        line=dict(color='red', width=3),
                        text=['', f'X: {bbox["size"][0]:.2f}'],
                        textposition='top center',
                        name='X dimension'
                    ))
                
                fig.update_layout(
                    scene=dict(
                        aspectmode='data',
                        camera=dict(eye=dict(x=1.5, y=1.5, z=1.0)),
                        bgcolor='#1a1a2e'
                    ),
                    paper_bgcolor='#1a1a2e',
                    margin=dict(l=0, r=0, t=0, b=0),
                    height=500,
                    clickmode='event+select'
                )
                
                # Use plotly_events for click detection
                from streamlit_plotly_events import plotly_events
                
                selected_points = plotly_events(
                    fig,
                    click_event=True,
                    hover_event=False,
                    select_event=False,
                    override_height=500,
                    key="mesh_click"
                )
                
                # Handle click events
                if selected_points:
                    st.success(f"識 繧ｯ繝ｪ繝・け讀懷・: {selected_points}")
                    # Get clicked point info
                    for pt in selected_points:
                        if 'pointNumber' in pt:
                            face_idx = pt['pointNumber'] // 3  # Approximate face index
                            st.write(f"驕ｸ謚槭＆繧後◆髱｢: 邏・{face_idx}")
                            
                            # Add dimension from clicked face
                            if st.session_state.extracted_dims and face_idx < len(st.session_state.extracted_dims):
                                dim = st.session_state.extracted_dims[face_idx]
                                if st.button(f"筐・{dim['name']} 繧偵メ繧ｧ繝ｼ繝ｳ縺ｫ霑ｽ蜉", key=f"add_click_{face_idx}"):
                                    st.session_state.tol_dimensions.append({
                                        "name": dim['name'],
                                        "nominal": dim['nominal'],
                                        "tolerance": dim.get('tolerance', 0.1),
                                        "direction": "+"
                                    })
                                    st.rerun()
                
                # Face selection UI
                st.markdown("---")
                st.subheader("Face Selection")
                if st.session_state.extracted_dims:
                    # Create face selector with colors
                    face_options = {f"{chr(65+i)}: {d['name']} ({d['nominal']:.2f}mm)": i 
                                   for i, d in enumerate(st.session_state.extracted_dims[:20])}
                    
                    selected_face = st.selectbox("Select dimension", list(face_options.keys()), key="face_select")
                    
                    col_sel1, col_sel2 = st.columns(2)
                    with col_sel1:
                        direction = st.radio("Direction", ["+", "-"], horizontal=True, key="dir_select")
                    with col_sel2:
                        if st.button("笵難ｸ・繝√ぉ繝ｼ繝ｳ縺ｫ霑ｽ蜉", type="primary", use_container_width=True):
                            idx = face_options[selected_face]
                            dim = st.session_state.extracted_dims[idx]
                            st.session_state.tol_dimensions.append({
                                "name": dim['name'],
                                "nominal": dim['nominal'],
                                "tolerance": dim.get('tolerance', 0.1),
                                "direction": direction
                            })
                            st.success(f"笨・{dim['name']} 繧定ｿｽ蜉縺励∪縺励◆")
                            st.rerun()
                else:
                    st.caption("棟 繝輔ぃ繧､繝ｫ繧偵い繝・・繝ｭ繝ｼ繝峨＠縺ｦ蟇ｸ豕輔ｒ謚ｽ蜃ｺ縺励※縺上□縺輔＞")
            else:
                st.info("刀 蟾ｦ蛛ｴ縺九ｉSTEP/STL繝輔ぃ繧､繝ｫ繧偵い繝・・繝ｭ繝ｼ繝峨＠縺ｦ3D繝｢繝・Ν繧定｡ｨ遉ｺ")
                
                # Show placeholder 3D
                import plotly.graph_objects as go
                import numpy as np
                
                # Demo box
                fig = go.Figure(data=[
                    go.Mesh3d(
                        x=[0, 0, 1, 1, 0, 0, 1, 1],
                        y=[0, 1, 1, 0, 0, 1, 1, 0],
                        z=[0, 0, 0, 0, 1, 1, 1, 1],
                        i=[0, 0, 1, 1, 4, 4, 0, 2, 1, 5, 0, 4],
                        j=[1, 2, 2, 3, 5, 6, 4, 3, 5, 6, 1, 5],
                        k=[2, 3, 6, 7, 6, 7, 1, 7, 6, 2, 4, 1],
                        color='#667eea',
                        opacity=0.5,
                        flatshading=True
                    )
                ])
                fig.update_layout(
                    scene=dict(aspectmode='cube', bgcolor='#1a1a2e'),
                    paper_bgcolor='#1a1a2e',
                    margin=dict(l=0, r=0, t=30, b=0),
                    height=400,
                    title="Demo: placeholder model"
                )
                st.plotly_chart(fig, use_container_width=True)
    
    # === Tab 2: Tolerance Chain ===
    with tab_dims:
        col_add, col_chain = st.columns([1, 2])
        
        with col_add:
            st.subheader("筐・蟇ｸ豕輔ｒ霑ｽ蜉")
            
            with st.form("add_dim_form"):
                dim_name = st.text_input("蟇ｸ豕募錐", f"DIM_{len(st.session_state.tol_dimensions)+1}")
                dim_nominal = st.number_input("蜈ｬ遘ｰ蛟､ (mm)", value=10.0, step=0.1)
                dim_tol = st.number_input("蜈ｬ蟾ｮ (ﾂｱmm)", value=0.1, step=0.01, min_value=0.001)
                dim_direction = st.radio("Direction", ["+", "-"], horizontal=True,
                                        help="+ = positive chain direction, - = negative chain direction")
                
                if st.form_submit_button("筐・繝√ぉ繝ｼ繝ｳ縺ｫ霑ｽ蜉", use_container_width=True):
                    st.session_state.tol_dimensions.append({
                        "name": dim_name,
                        "nominal": dim_nominal,
                        "tolerance": dim_tol,
                        "direction": dim_direction
                    })
                    st.rerun()
            
            if st.button("卵・・繝√ぉ繝ｼ繝ｳ繧偵け繝ｪ繧｢", use_container_width=True):
                st.session_state.tol_dimensions = []
                st.session_state.tol_result = None
                st.rerun()
        
        with col_chain:
            st.subheader("笵難ｸ・蜈ｬ蟾ｮ繝√ぉ繝ｼ繝ｳ")
            
            if st.session_state.tol_dimensions:
                dims_df = pd.DataFrame(st.session_state.tol_dimensions)
                edited_dims = st.data_editor(
                    dims_df,
                    use_container_width=True,
                    num_rows="dynamic",
                    column_config={
                        "name": st.column_config.TextColumn("蟇ｸ豕募錐"),
                        "nominal": st.column_config.NumberColumn("蜈ｬ遘ｰ蛟､ (mm)", format="%.4f"),
                        "tolerance": st.column_config.NumberColumn("蜈ｬ蟾ｮ (ﾂｱmm)", format="%.4f"),
                        "direction": st.column_config.SelectboxColumn("Direction", options=["+", "-"])
                    }
                )
                st.session_state.tol_dimensions = edited_dims.to_dict('records')
                
                # Chain summary
                chain_nominal = sum(
                    d["nominal"] * (1 if d["direction"] == "+" else -1)
                    for d in st.session_state.tol_dimensions
                )
                
                col_m1, col_m2 = st.columns(2)
                with col_m1:
                    st.metric("繝√ぉ繝ｼ繝ｳ蜈ｬ遘ｰ蛟､", f"{chain_nominal:.4f} mm")
                with col_m2:
                    st.metric("蟇ｸ豕墓焚", len(st.session_state.tol_dimensions))
                
                if st.button("Run Analysis", use_container_width=True, type="primary"):
                    import math
                    import random
                    
                    dims = st.session_state.tol_dimensions
                    
                    # Worst Case
                    wc_upper = sum(d["tolerance"] for d in dims)
                    
                    # RSS (3ﾏ・
                    rss = math.sqrt(sum(d["tolerance"]**2 for d in dims))
                    
                    # Monte Carlo
                    mc_samples = 10000
                    mc_results = []
                    for _ in range(mc_samples):
                        sample = sum(
                            (d["nominal"] + random.gauss(0, d["tolerance"]/3)) * (1 if d["direction"] == "+" else -1)
                            for d in dims
                        )
                        mc_results.append(sample)
                    
                    mc_mean = sum(mc_results) / len(mc_results)
                    mc_std = math.sqrt(sum((x - mc_mean)**2 for x in mc_results) / len(mc_results))
                    
                    # Sensitivity
                    total_var = sum(d["tolerance"]**2 for d in dims)
                    sensitivities = {d["name"]: (d["tolerance"]**2 / total_var * 100) if total_var > 0 else 0 for d in dims}
                    
                    st.session_state.tol_result = {
                        "nominal": chain_nominal,
                        "wc_upper": wc_upper,
                        "rss": rss,
                        "mc_mean": mc_mean,
                        "mc_std": mc_std,
                        "sensitivities": sensitivities,
                        "mc_histogram": mc_results
                    }
                    st.rerun()
            else:
                st.info("盗 蟾ｦ蛛ｴ縺九ｉ蟇ｸ豕輔ｒ霑ｽ蜉縺吶ｋ縺九・D繝｢繝・Ν繧ｿ繝悶〒蟇ｸ豕輔ｒ驕ｸ謚槭＠縺ｦ縺上□縺輔＞")
    
    # === Tab 3: Results ===
    with tab_result:
        if st.session_state.tol_result:
            r = st.session_state.tol_result
            
            # Summary cards
            st.subheader("投 隗｣譫舌し繝槭Μ")
            col_r1, col_r2, col_r3, col_r4 = st.columns(4)
            with col_r1:
                st.metric("蜈ｬ遘ｰ蛟､", f"{r['nominal']:.4f} mm")
            with col_r2:
                st.metric("Worst Case", f"ﾂｱ{r['wc_upper']:.4f} mm", delta_color="off")
            with col_r3:
                st.metric("RSS (3ﾏ・", f"ﾂｱ{r['rss']:.4f} mm", delta_color="off")
            with col_r4:
                st.metric("Monte Carlo σ", f"{r['mc_std']:.4f} mm", delta_color="off")
            
            # Results comparison table
            st.markdown("---")
            results_df = pd.DataFrame([
                {"Method": "Worst Case", "Upper": r['nominal'] + r['wc_upper'], "Lower": r['nominal'] - r['wc_upper'], "Span": r['wc_upper'] * 2},
                {"Method": "RSS (3σ)", "Upper": r['nominal'] + r['rss'], "Lower": r['nominal'] - r['rss'], "Span": r['rss'] * 2},
                {"Method": "Monte Carlo (3σ)", "Upper": r['mc_mean'] + 3*r['mc_std'], "Lower": r['mc_mean'] - 3*r['mc_std'], "Span": 6*r['mc_std']},
            ])
            st.dataframe(results_df, use_container_width=True, hide_index=True)
            
            # Sensitivity chart
            col_sens, col_hist = st.columns(2)
            
            with col_sens:
                st.subheader("識 諢溷ｺｦ蛻・梵")
                sens_df = pd.DataFrame([{"Dimension": k, "Sensitivity": v} for k, v in r['sensitivities'].items()])
                sens_df = sens_df.sort_values("Sensitivity", ascending=True)
                
                import plotly.express as px
                fig_sens = px.bar(sens_df, x="Sensitivity", y="Dimension", orientation='h',
                                 color="Sensitivity", color_continuous_scale="Blues")
                fig_sens.update_layout(height=300, showlegend=False)
                st.plotly_chart(fig_sens, use_container_width=True)
            
            with col_hist:
                st.subheader("Monte Carlo Histogram")
                import plotly.express as px
                
                fig_hist = px.histogram(x=r['mc_histogram'], nbins=50)
                fig_hist.add_vline(x=r['mc_mean'], line_dash="dash", line_color="red", annotation_text="ﾎｼ")
                fig_hist.add_vline(x=r['mc_mean'] - 3*r['mc_std'], line_dash="dot", line_color="orange")
                fig_hist.add_vline(x=r['mc_mean'] + 3*r['mc_std'], line_dash="dot", line_color="orange")
                fig_hist.update_layout(height=300, xaxis_title="邏ｯ遨榊・蟾ｮ (mm)", yaxis_title="鬆ｻ蠎ｦ")
                st.plotly_chart(fig_hist, use_container_width=True)
            
            # Export buttons
            st.markdown("---")
            col_exp1, col_exp2 = st.columns(2)
            
            with col_exp1:
                report_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>蜈ｬ蟾ｮ隗｣譫舌Ξ繝昴・繝・/title>
<style>body{{font-family:sans-serif;padding:20px;background:#1a1a2e;color:#eee}}
table{{border-collapse:collapse;width:100%}}th,td{{border:1px solid #444;padding:8px}}
th{{background:#667eea}}.metric{{background:#2d2d44;padding:15px;border-radius:8px;margin:5px}}</style></head>
<body><h1>投 蜈ｬ蟾ｮ隗｣譫舌Ξ繝昴・繝・/h1>
<p>逕滓・: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
<div style="display:flex;gap:10px;">
<div class="metric"><h3>蜈ｬ遘ｰ蛟､</h3><p>{r['nominal']:.4f} mm</p></div>
<div class="metric"><h3>Worst Case</h3><p>ﾂｱ{r['wc_upper']:.4f} mm</p></div>
<div class="metric"><h3>RSS</h3><p>ﾂｱ{r['rss']:.4f} mm</p></div>
<div class="metric"><h3>MC ﾏ・/h3><p>{r['mc_std']:.4f} mm</p></div>
</div>
<h2>蟇ｸ豕穂ｸ隕ｧ</h2><table><tr><th>蜷咲ｧｰ</th><th>蜈ｬ遘ｰ蛟､</th><th>蜈ｬ蟾ｮ</th><th>譁ｹ蜷・/th><th>蟇・ｸ主ｺｦ</th></tr>
{"".join(f'<tr><td>{d["name"]}</td><td>{d["nominal"]:.4f}</td><td>ﾂｱ{d["tolerance"]:.4f}</td><td>{d["direction"]}</td><td>{r["sensitivities"].get(d["name"],0):.1f}%</td></tr>' for d in st.session_state.tol_dimensions)}
</table></body></html>"""
                st.download_button("Download HTML Report", report_html, "tolerance_report.html", "text/html", use_container_width=True)
            
            with col_exp2:
                csv_data = "蜷咲ｧｰ,蜈ｬ遘ｰ蛟､,蜈ｬ蟾ｮ,譁ｹ蜷・蟇・ｸ主ｺｦ\n" + "\n".join(
                    f'{d["name"]},{d["nominal"]},{d["tolerance"]},{d["direction"]},{r["sensitivities"].get(d["name"],0):.1f}'
                    for d in st.session_state.tol_dimensions
                )
                st.download_button("Download CSV", csv_data, "tolerance_data.csv", "text/csv", use_container_width=True)
        else:
            st.info("Add dimensions in the tolerance table to run the analysis.")

elif page == "Kindle Manuscript":
    st.header("答 譖ｸ邀榊次遞ｿ逕滓・ (Kindle Unlimited)")
    st.markdown("""
    **菴ｿ縺・婿:**
    1. `/consume/Kindle/` 縺ｫ繝励Ο繧ｸ繧ｧ繧ｯ繝医ヵ繧ｩ繝ｫ繝繧剃ｽ懈・・井ｾ・ `FEM_Impact`・・
    2. 雉・侭・・DF, PPTX, TXT, DOCX・峨ｒ繝輔か繝ｫ繝縺ｫ驟咲ｽｮ
    3. IMPACT FEM繧定ｵｷ蜍輔＠縺ｦ繧ｹ繧ｯ繝ｪ繝ｼ繝ｳ繧ｷ繝ｧ繝・ヨ繧貞叙蠕暦ｼ医が繝励す繝ｧ繝ｳ・・
    4. AI縺瑚ｳ・侭縺ｨ逕ｻ蜒上ｒ隱ｭ縺ｿ霎ｼ縺ｿ縲∝次遞ｿ繧堤函謌・
    
    ---
    """)
    
    # List available Kindle project folders
    kindle_projects = []
    if os.path.exists(KINDLE_DIR):
        kindle_projects = [d for d in os.listdir(KINDLE_DIR) if os.path.isdir(os.path.join(KINDLE_DIR, d))]
    
    if not kindle_projects:
        st.warning(f"No Kindle projects were found. Place project folders under `{KINDLE_DIR}/`.")
    else:
        selected_project = st.selectbox("Select project", kindle_projects)
        project_path = os.path.join(KINDLE_DIR, selected_project)
        images_dir = os.path.join(project_path, "images")
        os.makedirs(images_dir, exist_ok=True)
        
        # List files in selected project (Recursively)
        project_files = []
        for root, dirs, files in os.walk(project_path):
            for f in files:
                # Exclude hidden files or system files if needed
                if not f.startswith('.'):
                    # Store relative path for cleaner display
                    rel_path = os.path.relpath(os.path.join(root, f), project_path)
                    project_files.append(rel_path)
        
        st.write(f"**雉・侭繝輔ぃ繧､繝ｫ ({len(project_files)}莉ｶ):**")
        for f in project_files[:10]:  # Show first 10
            fpath = os.path.join(project_path, f)
            size_kb = os.path.getsize(fpath) / 1024
            st.write(f"- {f} ({size_kb:.1f} KB)")
        if len(project_files) > 10:
            st.write(f"... 莉・{len(project_files) - 10} 莉ｶ")
        
        # IMPACT Control Section
        st.markdown("---")
        st.subheader("萄 IMPACT FEM 繧ｹ繧ｯ繝ｪ繝ｼ繝ｳ繧ｷ繝ｧ繝・ヨ")
        
        impact_col1, impact_col2, impact_col3 = st.columns(3)
        
        with impact_col1:
            if st.button("Start IMPACT", use_container_width=True):
                import subprocess
                try:
                    result = subprocess.run(
                        ["bash", "/work/scripts/impact_vnc.sh", "start"],
                        capture_output=True, text=True, timeout=30
                    )
                    if result.returncode == 0:
                        st.success("笨・IMPACT襍ｷ蜍募ｮ御ｺ・")
                        st.info("倹 http://localhost:6080/vnc.html 縺ｧGUI謫堺ｽ懷庄閭ｽ")
                    else:
                        st.error(f"繧ｨ繝ｩ繝ｼ: {result.stderr or result.stdout}")
                except Exception as e:
                    st.error(f"襍ｷ蜍募､ｱ謨・ {e}")
        
        with impact_col2:
            screenshot_name = st.text_input("Screenshot name", f"screen_{datetime.datetime.now().strftime('%H%M%S')}")
            if st.button("Capture Screenshot", use_container_width=True):
                import subprocess
                screenshot_path = os.path.join(images_dir, f"{screenshot_name}.png")
                try:
                    result = subprocess.run(
                        ["bash", "/work/scripts/impact_vnc.sh", "screenshot", screenshot_path],
                        capture_output=True, text=True, timeout=10
                    )
                    if result.returncode == 0 and os.path.exists(screenshot_path):
                        st.success(f"笨・菫晏ｭ・ {screenshot_path}")
                        st.image(screenshot_path, caption=screenshot_name, width=300)
                    else:
                        st.error(f"繧ｨ繝ｩ繝ｼ: {result.stderr or result.stdout}")
                except Exception as e:
                    st.error(f"蜿門ｾ怜､ｱ謨・ {e}")
        
        with impact_col3:
            if st.button("尅 IMPACT蛛懈ｭ｢", use_container_width=True):
                import subprocess
                try:
                    subprocess.run(["bash", "/work/scripts/impact_vnc.sh", "stop"], timeout=10)
                    st.success("IMPACT stopped.")
                except Exception as e:
                    st.warning(f"蛛懈ｭ｢荳ｭ縺ｫ繧ｨ繝ｩ繝ｼ: {e}")
        
        # Show existing images
        existing_images = [f for f in os.listdir(images_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))] if os.path.exists(images_dir) else []
        if existing_images:
            with st.expander(f"胴 菫晏ｭ俶ｸ医∩逕ｻ蜒・({len(existing_images)}譫・"):
                img_cols = st.columns(3)
                for i, img in enumerate(existing_images[:9]):
                    with img_cols[i % 3]:
                        st.image(os.path.join(images_dir, img), caption=img, width=150)
        
        # ---------------------------------------------------------
        # 竭｢ 繧ｷ繝溘Η繝ｬ繝ｼ繧ｷ繝ｧ繝ｳ螳溯｡・(P018)
        # ---------------------------------------------------------
        st.markdown("---")
        st.subheader("竭｢ Impact FEM 繧ｷ繝溘Η繝ｬ繝ｼ繧ｷ繝ｧ繝ｳ螳溯｡・(P018)")
        
        # Filter .in files
        in_files = [f for f in project_files if f.lower().endswith(".in")]
        
        if not in_files:
            st.info("No `.in` input file was found.")
        else:
            selected_in_file = st.selectbox("蟇ｾ雎｡縺ｮ蜈･蜉帙ヵ繧｡繧､繝ｫ (.in)", in_files)
            target_in_path = os.path.join(project_path, selected_in_file)
            
            c1, c2 = st.columns(2)
            
            with c1:
                if st.button("噫 繧ｷ繝溘Η繝ｬ繝ｼ繧ｷ繝ｧ繝ｳ螳溯｡・(Background)", use_container_width=True):
                    # Command construction for Container Environment
                    # Classpath assumes /opt/impact structure
                    classpath = "/opt/impact/bin:/opt/impact/lib/*"
                    java_cmd = [
                        "java", 
                        "-Xmx2048m", 
                        "-cp", classpath,
                        "run.Impact", 
                        target_in_path
                    ]
                    
                    st.write(f"Executing: `{' '.join(java_cmd)}`")
                    
                    try:
                        # Run in background (nohup style) or blocking? 
                        # For simple UX, let's do blocking with a spinner for now, 
                        # or use subprocess.Popen for background if long running.
                        # Given Streamlit's nature, blocking is easier to show logs, 
                        # but for long sims, background is better.
                        # Let's try blocking for immediate feedback on simple models like Bullet_AKM.
                        
                        with st.spinner("繧ｷ繝溘Η繝ｬ繝ｼ繧ｷ繝ｧ繝ｳ螳溯｡御ｸｭ... (繝ｭ繧ｰ縺ｯ荳九↓陦ｨ遉ｺ縺輔ｌ縺ｾ縺・"):
                            process = subprocess.Popen(
                                java_cmd, 
                                stdout=subprocess.PIPE, 
                                stderr=subprocess.PIPE, 
                                text=True, 
                                cwd=os.path.dirname(target_in_path) # Run in file's dir
                            )
                            stdout, stderr = process.communicate()
                            
                            if process.returncode == 0:
                                st.success("Simulation completed.")
                            else:
                                st.error(f"笶・繧ｨ繝ｩ繝ｼ逋ｺ逕・(Exit Code: {process.returncode})")
                            
                            with st.expander("螳溯｡後Ο繧ｰ (STDOUT)", expanded=True):
                                st.code(stdout)
                            if stderr:
                                with st.expander("繧ｨ繝ｩ繝ｼ繝ｭ繧ｰ (STDERR)", expanded=True):
                                    st.code(stderr)
                                    
                    except Exception as e:
                        st.error(f"Execution Error: {e}")

            with c2:
                if st.button("祷 縺薙・隗｣譫舌・隗｣隱ｬ逕滓・ (AI)", use_container_width=True):
                    with st.spinner("隗｣譫仙・螳ｹ繧貞・譫蝉ｸｭ..."):
                         # Read the .in file
                        try:
                            with open(target_in_path, "r", encoding="utf-8", errors="ignore") as f:
                                in_content = f.read()
                            
                            prompt = f"""
縺ゅ↑縺溘・CAE隗｣譫舌・蟆る摩螳ｶ縺ｧ縺吶ゆｻ･荳九・Impact FEM蜈･蜉帙ヵ繧｡繧､繝ｫ(`{selected_in_file}`)繧貞・譫舌＠縲・
縺薙・繧ｷ繝溘Η繝ｬ繝ｼ繧ｷ繝ｧ繝ｳ縺後御ｽ輔ｒ縲阪後←縺・＞縺・擅莉ｶ縺ｧ縲崎ｧ｣譫舌＠繧医≧縺ｨ縺励※縺・ｋ縺ｮ縺九・
荳闊ｬ縺ｮ繧ｨ繝ｳ繧ｸ繝九い縺ｫ繧ゅｏ縺九ｋ繧医≧縺ｫ隗｣隱ｬ繝ｬ繝昴・繝医ｒ菴懈・縺励※縺上□縺輔＞縲・

# 鬆・岼
1. **隗｣譫舌・逶ｮ逧・*: 菴輔′菴輔↓陦晉ｪ√☆繧九・縺九√↑縺ｩ
2. **繝｢繝・Ν讎りｦ・*: 繝弱・繝画焚縲∬ｦ∫ｴ繧ｿ繧､繝励∵攝譁呻ｼ・aterials・・
3. **蠅・阜譚｡莉ｶ (Constraints)**: 縺ｩ縺薙′蝗ｺ螳壹＆繧後※縺・ｋ縺九∝・騾溷ｺｦ縺ｯ縺・￥繧峨°
4. **譛溷ｾ・＆繧後ｋ邨先棡**: 縺ｩ縺ｮ繧医≧縺ｪ迚ｩ逅・樟雎｡・郁ｲｫ騾壹∬ｷｳ縺ｭ霑斐ｊ縲∝､牙ｽ｢・峨′隕九ｉ繧後ｋ縺ｯ縺壹°

# 蜈･蜉帙ヵ繧｡繧､繝ｫ蜀・ｮｹ
{in_content[:20000]}
"""
                            explanation = ask_ai(prompt)
                            st.markdown(explanation)
                        except Exception as e:
                            st.error(f"Analysis Error: {e}")

        st.markdown("---")
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("竭 雉・侭隱ｭ縺ｿ霎ｼ縺ｿ")
            include_code = st.checkbox("繧ｽ繝ｼ繧ｹ繧ｳ繝ｼ繝・.py, .java, .md遲・繧ょ性繧√ｋ", value=True)
            if st.button("Analyze Materials", use_container_width=True):
                combined_text = ""
                for f in project_files:
                    fpath = os.path.join(project_path, f)
                    ext = f.lower().split('.')[-1]
                    try:
                        # Text / Code Files
                        if ext in ["txt", "sh", "py", "java", "c", "cpp", "h", "md", "json", "yml", "yaml", "bat", "ps1", "properties", "in"]:
                            if include_code or ext == "txt":
                                with open(fpath, "r", encoding="utf-8", errors="ignore") as tf:
                                    combined_text += f"\n\n=== FILE: {f} ===\n" + tf.read()
                        elif ext == "pdf":
                            reader = pypdf.PdfReader(fpath)
                            combined_text += f"\n\n=== {f} ===\n" + "\n".join([p.extract_text() or "" for p in reader.pages])
                        elif ext == "docx":
                            doc = docx.Document(fpath)
                            combined_text += f"\n\n=== {f} ===\n" + "\n".join([p.text for p in doc.paragraphs])
                        elif ext == "pptx":
                            from pptx import Presentation
                            prs = Presentation(fpath)
                            slides_text = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        slides_text.append(shape.text)
                            combined_text += f"\n\n=== {f} ===\n" + "\n".join(slides_text)
                    except Exception as e:
                        st.warning(f"{f}: 隱ｭ縺ｿ霎ｼ縺ｿ繧ｨ繝ｩ繝ｼ ({e})")
                
                st.session_state.kindle_content = combined_text
                st.session_state.kindle_images = existing_images
                st.success(f"Loaded {len(combined_text)} characters and {len(existing_images)} images.")
                with st.expander("隱ｭ縺ｿ霎ｼ縺ｿ蜀・ｮｹ繝励Ξ繝薙Η繝ｼ"):
                    st.text(combined_text[:3000] + "..." if len(combined_text) > 3000 else combined_text)
        
        with col2:
            st.subheader("竭｡ 蜴溽ｨｿ逕滓・")
            book_title = st.text_input("譖ｸ邀阪ち繧､繝医Ν", f"{selected_project}蜈･髢")
            target_pages = st.number_input("逶ｮ讓吶・繝ｼ繧ｸ謨ｰ", min_value=10, max_value=200, value=50)
            
            if st.button("笨ｨ AI縺ｧ蜴溽ｨｿ逕滓・", use_container_width=True):
                if "kindle_content" not in st.session_state or not st.session_state.kindle_content:
                    st.error("蜈医↓縲瑚ｳ・侭繧定ｧ｣譫舌阪ｒ螳溯｡後＠縺ｦ縺上□縺輔＞")
                else:
                    # Build image references
                    image_refs = ""
                    if existing_images:
                        image_refs = "\n\n蛻ｩ逕ｨ蜿ｯ閭ｽ縺ｪ逕ｻ蜒・\n" + "\n".join([f"- {img}" for img in existing_images])
                    
                    prompt = f"""
Create a Kindle-ready Markdown manuscript.

Title: {book_title}
Target pages: {target_pages}
Target length guideline: about {target_pages * 400} Japanese characters.

Requirements:
- Organize the material into a practical, readable book structure.
- Use the supplied technical materials and screenshots as source material.
- Explain engineering concepts clearly for readers.
- Return Markdown only.

Available images:
{image_refs}

Source materials:
{st.session_state.kindle_content[:150000]}
"""
                    with st.spinner("蜴溽ｨｿ逕滓・荳ｭ... (謨ｰ蛻・°縺九ｋ蝣ｴ蜷医′縺ゅｊ縺ｾ縺・"):
                        result = ask_ai(prompt)
                    
                    st.session_state.kindle_manuscript = result
                    st.success("Manuscript generation completed.")
        
        # Display manuscript
        if "kindle_manuscript" in st.session_state and st.session_state.kindle_manuscript:
            st.markdown("---")
            st.subheader("当 逕滓・縺輔ｌ縺溷次遞ｿ")
            st.markdown(st.session_state.kindle_manuscript)
            
            # Save button
            fn = f"{selected_project}_蜴溽ｨｿ_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.md"
            if st.download_button("Download Markdown", st.session_state.kindle_manuscript, file_name=fn):
                st.success(f"菫晏ｭ伜ｮ御ｺ・ {fn}")

        # Check for existing reports
        st.markdown("---")
        st.subheader("唐 菫晏ｭ俶ｸ医∩蜴溽ｨｿ")
        report_files = [f for f in os.listdir(WORK_DIR) if f.startswith(selected_project) and f.endswith(".md")] if os.path.exists(WORK_DIR) else []
        
        if report_files:
            for rf in report_files:
                rf_path = os.path.join(WORK_DIR, rf)
                with open(rf_path, "r", encoding="utf-8") as f:
                    content = f.read()
                
                col_d1, col_d2 = st.columns([3, 1])
                with col_d1:
                    st.text(f"塘 {rf} ({os.path.getsize(rf_path)/1024:.1f} KB)")
                with col_d2:
                    st.download_button("Download", content, file_name=rf, key=f"dl_{rf}")
        else:
            st.info("No saved manuscript reports were found.")

# -------------------------
# P016: Email Reporting
# -------------------------
elif page == "Email Daily Report (P016)":
    st.header("透 Email Daily Report (P016)")
    st.info("P016 bundles requests, QIF items, and meeting notes into a daily report.")

    col1, col2 = st.columns([1, 2])
    
    with col1:
        # Date Filter
        today = datetime.date.today()
        start_date = st.date_input("髢句ｧ区律", today)
        end_date = st.date_input("邨ゆｺ・律", today)
        
        if st.button("Generate Report", use_container_width=True):
            with st.spinner("繝｡繝ｼ繝ｫ隗｣譫蝉ｸｭ... (謨ｰ蛻・°縺九ｊ縺ｾ縺・"):
                # Copy script if not present (safety check)
                if not os.path.exists("/app/generate_email_report.py"):
                    try:
                        import shutil
                        shutil.copy("/work/scripts/generate_email_report.py", "/app/generate_email_report.py")
                    except:
                        pass
                
                # Execute with Date Arguments
                cmd = [
                    "python3", "/work/scripts/generate_email_report.py",
                    "--start", str(start_date),
                    "--end", str(end_date)
                ]
                try:
                    result = subprocess.run(cmd, capture_output=True, text=True, check=True)
                    st.success("笨・繝ｬ繝昴・繝育函謌仙ｮ御ｺ・")
                    st.text_area("Log Output", result.stdout, height=200)
                except subprocess.CalledProcessError as e:
                    st.error(f"Error: {e}")
                    st.text_area("Error Output", e.stderr, height=200)

    # List Reports
    st.markdown("---")
    st.subheader("Available Reports")
    
    if os.path.exists(WORK_DIR):
        reports = [f for f in os.listdir(WORK_DIR) if f.startswith("Email_Report_") and f.endswith(".md")]
        reports.sort(reverse=True)
        
        if reports:
            selected_report = st.selectbox("Select report", reports)
            
            if selected_report:
                rpath = os.path.join(WORK_DIR, selected_report)
                with open(rpath, "r", encoding="utf-8") as f:
                    content = f.read()
                
                st.markdown(content)
                st.download_button("Download", content, file_name=selected_report)
        else:
            st.info("No reports are available yet.")

