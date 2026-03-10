#!/usr/bin/env python3
"""
ingest_pan_to_qdrant.py
ミツイ精密 ＰＡN・異常連絡書 フォルダの文書を Qdrant universal_knowledge に埋め込む。

対象:
  /home/node/clawd/paperless_consume/ミツイ精密/ＰＡN・異常連絡書/**
  - .xlsx / .xls : 不適合管理台帳・なぜなぜ分析・異常連絡書
  - .pptx / .ppt : プレゼンテーション資料
  - .docx / .doc : Word文書

スキップ: .db .jpg .png .gif .tif .lnk .zip .mp4 .xlk

実行:
  docker exec clawstack-unified-clawdbot-gateway-1 \
    python3 /home/node/clawd/ingest_pan_to_qdrant.py
"""

import os
import sys
import json
import hashlib
import requests
import traceback
from pathlib import Path
from datetime import datetime

# ── Configuration ────────────────────────────────────────────────────────────────
PAN_ROOT     = "/home/node/clawd/paperless_consume/ミツイ精密/ＰＡN・異常連絡書"
STATE_FILE   = "/home/node/clawd/ingest_pan_state.json"
LOG_FILE     = "/home/node/clawd/ingest_pan.log"

INFINITY_URL = "http://infinity:7997/embeddings"
QDRANT_URL   = "http://qdrant:6333"
COLLECTION   = "universal_knowledge"
EMBED_MODEL  = "mixedbread-ai/mxbai-embed-large-v1"
EMBED_DIM    = 1024

CHUNK_SIZE     = 800
ROWS_PER_CHUNK = 6    # Excel: 何行まとめてチャンクにするか
EMBED_BATCH    = 6
QDRANT_BATCH   = 32

TARGET_EXTS = {".xlsx", ".xls", ".pptx", ".ppt", ".docx", ".doc"}
SKIP_EXTS   = {".db", ".jpg", ".jpeg", ".png", ".gif", ".tif", ".tiff",
               ".lnk", ".zip", ".mp4", ".xlk", ".pyc"}

# ── Logging ──────────────────────────────────────────────────────────────────────
def log(msg, level="INFO"):
    ts   = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    line = f"[{ts}] [{level}] {msg}"
    if sys.stdout.isatty():
        print(line, flush=True)
    else:
        try:
            with open(LOG_FILE, "a", encoding="utf-8") as f:
                f.write(line + "\n")
        except Exception:
            pass

# ── State ─────────────────────────────────────────────────────────────────────────
def load_state():
    if os.path.exists(STATE_FILE):
        try:
            with open(STATE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"processed": {}}

def save_state(state):
    with open(STATE_FILE, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

# ── ファイルパスからメタデータ抽出 ───────────────────────────────────────────────
def path_meta(filepath):
    """フォルダ構造から年度・サブカテゴリを抽出。"""
    rel   = Path(filepath).relative_to(PAN_ROOT)
    parts = rel.parts
    year  = parts[0] if len(parts) > 1 and ("年" in parts[0] or parts[0].isdigit()) else ""
    subcat = parts[0] if not year and len(parts) > 1 else (parts[1] if len(parts) > 2 else "")
    return {"year": year, "subcat": subcat}

# ── XLSX/XLS ─────────────────────────────────────────────────────────────────────
def excel_sheets_to_text(filepath):
    """XLSX → openpyxl、XLS → xlrd でシートを読み込み、行テキストのリストを返す。
    戻り値: [(sheet_name, [row_text, ...])]
    """
    ext = Path(filepath).suffix.lower()
    result = []

    if ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        for sname in wb.sheetnames:
            ws = wb[sname]
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    rows_data.append(" | ".join(vals))
            if rows_data:
                result.append((sname, rows_data))
        wb.close()

    elif ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(str(filepath))
        for si in range(wb.nsheets):
            ws  = wb.sheet_by_index(si)
            rows_data = []
            for ri in range(ws.nrows):
                vals = [str(v).strip() for v in ws.row_values(ri) if str(v).strip() and str(v).strip() != "''"]
                if vals:
                    rows_data.append(" | ".join(vals))
            if rows_data:
                result.append((ws.name, rows_data))

    return result

def make_excel_chunks(filepath):
    stem = Path(filepath).stem
    meta = path_meta(filepath)
    rel  = str(Path(filepath).relative_to(PAN_ROOT))
    base = {
        "source":   f"pan/{rel}",
        "category": "品質異常・PAN",
        "filename": stem,
        "filetype": Path(filepath).suffix.lstrip("."),
        "year":     meta["year"],
        "subcat":   meta["subcat"],
    }

    try:
        sheets = excel_sheets_to_text(filepath)
    except Exception as e:
        log(f"  Excel read error {Path(filepath).name}: {e}", "WARN")
        return []

    chunks = []
    for sheet_name, rows in sheets:
        # 先頭行をヘッダーとして保持し、ROWS_PER_CHUNK 行ずつチャンク
        header = rows[0] if rows else ""
        data   = rows[1:] if len(rows) > 1 else []
        if not data:
            # ヘッダーのみ → 1チャンク
            text = f"[{stem}][{sheet_name}]\n{header}"
            pid  = int(hashlib.md5(f"pan::{filepath}::{sheet_name}::h".encode()).hexdigest()[:15], 16)
            chunks.append((text, {**base, "sheet": sheet_name, "chunk": 0, "content": text[:400]}, pid))
            continue

        for ci, start in enumerate(range(0, len(data), ROWS_PER_CHUNK)):
            batch = data[start:start + ROWS_PER_CHUNK]
            lines = [f"[{stem}][{sheet_name}]", f"ヘッダー: {header}"] + batch
            text  = "\n".join(lines)
            if len(text.strip()) < 20:
                continue
            pid = int(hashlib.md5(f"pan::{filepath}::{sheet_name}::c{ci}".encode()).hexdigest()[:15], 16)
            chunks.append((text, {**base, "sheet": sheet_name, "chunk": ci, "content": text[:400]}, pid))

    return chunks

# ── PPTX / PPT ───────────────────────────────────────────────────────────────────
def make_pptx_chunks(filepath):
    from pptx import Presentation
    stem = Path(filepath).stem
    meta = path_meta(filepath)
    rel  = str(Path(filepath).relative_to(PAN_ROOT))
    base = {
        "source":   f"pan/{rel}",
        "category": "品質異常・PAN",
        "filename": stem,
        "filetype": "pptx",
        "year":     meta["year"],
        "subcat":   meta["subcat"],
    }
    try:
        prs    = Presentation(str(filepath))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [s.text.strip() for s in slide.shapes
                     if hasattr(s, "text") and s.text.strip()]
            if parts:
                slides.append(f"[スライド{i}]\n" + "\n".join(parts))
        text = "\n\n".join(slides)
    except Exception as e:
        log(f"  PPTX error {Path(filepath).name}: {e}", "WARN")
        return []

    if len(text.strip()) < 20:
        return []
    chunks = []
    for ci, start in enumerate(range(0, len(text), CHUNK_SIZE)):
        chunk = text[start:start + CHUNK_SIZE]
        if not chunk.strip():
            continue
        pid = int(hashlib.md5(f"pan::{filepath}::c{ci}".encode()).hexdigest()[:15], 16)
        chunks.append((chunk, {**base, "chunk": ci, "content": chunk[:400]}, pid))
    return chunks

# ── DOCX / DOC ───────────────────────────────────────────────────────────────────
def make_docx_chunks(filepath):
    from docx import Document
    stem = Path(filepath).stem
    meta = path_meta(filepath)
    rel  = str(Path(filepath).relative_to(PAN_ROOT))
    base = {
        "source":   f"pan/{rel}",
        "category": "品質異常・PAN",
        "filename": stem,
        "filetype": Path(filepath).suffix.lstrip("."),
        "year":     meta["year"],
        "subcat":   meta["subcat"],
    }
    try:
        doc  = Document(str(filepath))
        text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        log(f"  DOCX error {Path(filepath).name}: {e}", "WARN")
        return []

    if len(text.strip()) < 20:
        return []
    chunks = []
    for ci, start in enumerate(range(0, len(text), CHUNK_SIZE)):
        chunk = text[start:start + CHUNK_SIZE]
        if not chunk.strip():
            continue
        pid = int(hashlib.md5(f"pan::{filepath}::c{ci}".encode()).hexdigest()[:15], 16)
        chunks.append((chunk, {**base, "chunk": ci, "content": chunk[:400]}, pid))
    return chunks

# ── Embed & Upsert ───────────────────────────────────────────────────────────────
def batch_embed(texts):
    try:
        resp = requests.post(INFINITY_URL,
            json={"model": EMBED_MODEL, "input": texts}, timeout=120)
        resp.raise_for_status()
        return [d["embedding"] for d in resp.json()["data"]]
    except Exception as e:
        log(f"  Embed batch error: {e} — 1件ずつリトライ", "WARN")
        results = []
        for t in texts:
            try:
                r = requests.post(INFINITY_URL,
                    json={"model": EMBED_MODEL, "input": t}, timeout=120)
                r.raise_for_status()
                results.append(r.json()["data"][0]["embedding"])
            except Exception as e2:
                log(f"  Embed single error: {e2}", "WARN")
                results.append(None)
        return results

def batch_upsert(points):
    try:
        resp = requests.put(f"{QDRANT_URL}/collections/{COLLECTION}/points",
            json={"points": points}, timeout=30)
        resp.raise_for_status()
        return len(points)
    except Exception as e:
        log(f"  Qdrant error: {e}", "WARN")
        return 0

def flush_buffer(buf):
    if not buf:
        return
    vecs   = batch_embed([c[0] for c in buf])
    points = []
    for (text, payload, pid), vec in zip(buf, vecs):
        if vec and len(vec) == EMBED_DIM:
            points.append({"id": pid, "vector": vec,
                "payload": {**payload, "ingested_at": datetime.now().isoformat()}})
    for i in range(0, len(points), QDRANT_BATCH):
        batch_upsert(points[i:i + QDRANT_BATCH])
    buf.clear()

# ── Main ─────────────────────────────────────────────────────────────────────────
def main():
    log("=" * 60)
    log("  ingest_pan_to_qdrant.py 開始")
    log(f"  ソース: {PAN_ROOT}")
    log("=" * 60)

    state     = load_state()
    processed = state.get("processed", {})

    all_files = sorted(p for p in Path(PAN_ROOT).rglob("*")
                       if p.is_file() and p.suffix.lower() in TARGET_EXTS)
    pending   = [p for p in all_files if str(p) not in processed]

    by_ext = {}
    for p in pending:
        e = p.suffix.lower()
        by_ext[e] = by_ext.get(e, 0) + 1
    log(f"  未処理: {len(pending)}件 " + "  ".join(f"{e}={v}" for e,v in sorted(by_ext.items())))

    if not pending:
        log("  ✅ 新規ファイルなし。")
        return

    ok = skip = err = 0
    buf = []

    for i, path in enumerate(pending):
        ext = path.suffix.lower()
        try:
            if ext in (".xlsx", ".xls"):
                chunks = make_excel_chunks(str(path))
            elif ext in (".pptx", ".ppt"):
                chunks = make_pptx_chunks(str(path))
            elif ext in (".docx", ".doc"):
                chunks = make_docx_chunks(str(path))
            else:
                skip += 1
                continue

            if not chunks:
                processed[str(path)] = {"chunks": 0, "ts": datetime.now().isoformat()}
                skip += 1
                continue

            for item in chunks:
                buf.append(item)
            if len(buf) >= EMBED_BATCH:
                flush_buffer(buf)

            processed[str(path)] = {
                "chunks": len(chunks),
                "type":   ext.lstrip("."),
                "ts":     datetime.now().isoformat(),
            }
            ok += 1

        except Exception as e:
            log(f"  [{i+1}] ERROR {path.name}: {e}", "ERROR")
            log(traceback.format_exc(), "ERROR")
            processed[str(path)] = {"error": str(e), "ts": datetime.now().isoformat()}
            err += 1

        if (i + 1) % 30 == 0:
            flush_buffer(buf)
            state["processed"] = processed
            save_state(state)
            log(f"  [{i+1}/{len(pending)}]  ok={ok}  skip={skip}  err={err}")

    flush_buffer(buf)
    state["processed"] = processed
    save_state(state)
    log(f"  ✅ 完了: ok={ok}  skip={skip}  err={err}  合計={len(pending)}")

if __name__ == "__main__":
    main()
