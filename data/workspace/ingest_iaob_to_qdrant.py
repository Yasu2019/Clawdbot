#!/usr/bin/env python3
"""
ingest_iaob_to_qdrant.py
IAOB IATF 16949 監査ガイド (PPTX) を Qdrant universal_knowledge に埋め込む。

ソース: /home/node/clawd/paperless_consume/IAOB要約資料/**/*.pptx
  フォルダ構造: IAOB要約資料/{topic}/{filename}.pptx

パイプライン:
  PPTX → スライドテキスト抽出 → チャンク → Infinity embed → Qdrant upsert

実行:
  docker exec clawstack-unified-clawdbot-gateway-1 \
    python3 /home/node/clawd/ingest_iaob_to_qdrant.py
"""

import os
import sys
import json
import hashlib
import requests
import traceback
from pathlib import Path
from datetime import datetime

# ── Configuration ────────────────────────────────────────────────────────────────
IAOB_ROOT    = "/home/node/clawd/paperless_consume/IAOB要約資料"
STATE_FILE   = "/home/node/clawd/ingest_iaob_state.json"
LOG_FILE     = "/home/node/clawd/ingest_iaob.log"

INFINITY_URL = "http://infinity:7997/embeddings"
QDRANT_URL   = "http://qdrant:6333"
COLLECTION   = "universal_knowledge"
EMBED_MODEL  = "mixedbread-ai/mxbai-embed-large-v1"
EMBED_DIM    = 1024

CHUNK_SIZE   = 800
EMBED_BATCH  = 8
QDRANT_BATCH = 32

# ── Logging ──────────────────────────────────────────────────────────────────────
def log(msg, level="INFO"):
    ts   = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    line = f"[{ts}] [{level}] {msg}"
    if sys.stdout.isatty():
        print(line, flush=True)
    else:
        try:
            with open(LOG_FILE, "a", encoding="utf-8") as f:
                f.write(line + "\n")
        except Exception:
            pass

# ── State ─────────────────────────────────────────────────────────────────────────
def load_state():
    if os.path.exists(STATE_FILE):
        try:
            with open(STATE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"processed": {}}

def save_state(state):
    with open(STATE_FILE, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

# ── PPTX parsing ─────────────────────────────────────────────────────────────────
def extract_pptx_text(filepath):
    """PPTX からスライドテキストを抽出。スライド番号付きで結合。"""
    from pptx import Presentation
    prs = Presentation(str(filepath))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append(f"[スライド{i}]\n" + "\n".join(parts))
    return "\n\n".join(slides_text)

def parse_pptx(filepath):
    """PPTX をパースしてメタデータ + 本文を返す。"""
    try:
        rel      = Path(filepath).relative_to(IAOB_ROOT)
        topic    = rel.parts[0] if len(rel.parts) > 1 else "IAOB"
        filename = Path(filepath).stem

        text = extract_pptx_text(filepath)
        return {
            "filepath": filepath,
            "topic":    topic,
            "filename": filename,
            "text":     text,
        }
    except Exception as e:
        log(f"  Parse error {Path(filepath).name}: {e}", "WARN")
        return None

def make_chunks(meta):
    text = meta["text"].strip()
    if len(text) < 20:
        return []

    rel_path = str(Path(meta["filepath"]).relative_to(IAOB_ROOT))
    base_payload = {
        "source":   f"iaob/{rel_path}",
        "category": "IAOB_IATF16949",
        "topic":    meta["topic"],
        "filename": meta["filename"],
    }
    result = []
    for ci, start in enumerate(range(0, len(text), CHUNK_SIZE)):
        chunk = text[start:start + CHUNK_SIZE]
        if not chunk.strip():
            continue
        pid = int(hashlib.md5(f"iaob::{meta['filepath']}::c{ci}".encode()).hexdigest()[:15], 16)
        result.append((chunk, {**base_payload, "chunk": ci, "content": chunk}, pid))
    return result

# ── Infinity batch embed ─────────────────────────────────────────────────────────
def batch_embed(texts):
    try:
        resp = requests.post(
            INFINITY_URL,
            json={"model": EMBED_MODEL, "input": texts},
            timeout=120,
        )
        resp.raise_for_status()
        return [d["embedding"] for d in resp.json()["data"]]
    except Exception as e:
        log(f"  Embed batch error: {e} — 1件ずつリトライ", "WARN")
        results = []
        for t in texts:
            try:
                r = requests.post(INFINITY_URL,
                    json={"model": EMBED_MODEL, "input": t}, timeout=120)
                r.raise_for_status()
                results.append(r.json()["data"][0]["embedding"])
            except Exception as e2:
                log(f"  Embed single error: {e2}", "WARN")
                results.append(None)
        return results

# ── Qdrant batch upsert ──────────────────────────────────────────────────────────
def batch_upsert(points):
    try:
        resp = requests.put(
            f"{QDRANT_URL}/collections/{COLLECTION}/points",
            json={"points": points},
            timeout=30,
        )
        resp.raise_for_status()
        return len(points)
    except Exception as e:
        log(f"  Qdrant error: {e}", "WARN")
        return 0

# ── Main ─────────────────────────────────────────────────────────────────────────
def main():
    log("=" * 60)
    log("  ingest_iaob_to_qdrant.py 開始")
    log(f"  ソース: {IAOB_ROOT}/**/*.pptx")
    log("=" * 60)

    state     = load_state()
    processed = state.get("processed", {})

    all_files = sorted(Path(IAOB_ROOT).rglob("*.pptx"))
    pending   = [p for p in all_files if str(p) not in processed]

    log(f"  総ファイル数: {len(all_files)}  未処理: {len(pending)}")
    if not pending:
        log("  ✅ 新規ファイルなし。処理終了。")
        return

    ok = skip = err = 0
    chunk_buf = []  # (text, payload, pid, filepath)

    def flush():
        if not chunk_buf:
            return
        texts = [c[0] for c in chunk_buf]
        vecs  = batch_embed(texts)
        points = []
        for (text, payload, pid, fp), vec in zip(chunk_buf, vecs):
            if vec and len(vec) == EMBED_DIM:
                points.append({"id": pid, "vector": vec,
                    "payload": {**payload, "ingested_at": datetime.now().isoformat()}})
        for i in range(0, len(points), QDRANT_BATCH):
            batch_upsert(points[i:i + QDRANT_BATCH])
        chunk_buf.clear()

    for i, path in enumerate(pending):
        try:
            meta = parse_pptx(str(path))
            if not meta or not meta["text"].strip():
                processed[str(path)] = {"chunks": 0, "ts": datetime.now().isoformat()}
                skip += 1
                continue

            chunks = make_chunks(meta)
            if not chunks:
                processed[str(path)] = {"chunks": 0, "ts": datetime.now().isoformat()}
                skip += 1
                continue

            for (chunk_text, payload, pid) in chunks:
                chunk_buf.append((chunk_text, payload, pid, str(path)))

            if len(chunk_buf) >= EMBED_BATCH:
                flush()

            processed[str(path)] = {
                "topic":  meta["topic"][:60],
                "chunks": len(chunks),
                "ts":     datetime.now().isoformat(),
            }
            ok += 1

        except Exception as e:
            log(f"  [{i+1}] ERROR {path.name}: {e}", "ERROR")
            log(traceback.format_exc(), "ERROR")
            processed[str(path)] = {"error": str(e), "ts": datetime.now().isoformat()}
            err += 1

        if (i + 1) % 10 == 0:
            flush()
            state["processed"] = processed
            save_state(state)
            log(f"  [{i+1}/{len(pending)}]  ok={ok}  skip={skip}  err={err}")

    flush()
    state["processed"] = processed
    save_state(state)
    log(f"  ✅ 完了: ok={ok}  skip={skip}  err={err}  合計={len(pending)}")

if __name__ == "__main__":
    main()
