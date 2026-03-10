#!/usr/bin/env python3
"""
ingest_iatf_docs_to_qdrant.py
IATF_documents ディレクトリの PPTX / XLSX を Qdrant universal_knowledge に埋め込む。

ソース:
  /home/node/clawd/paperless_consume/IATF_documents/*.pptx  (97件: IATF規格箇条解説)
  /home/node/clawd/paperless_consume/IATF_documents/*.xlsx  (509件: FMEA/CP/会議事録等)

パイプライン:
  PPTX → python-pptx テキスト抽出
  XLSX → openpyxl テーブル行テキスト抽出
  → チャンク → Infinity embed → Qdrant upsert

実行:
  docker exec clawstack-unified-clawdbot-gateway-1 \
    python3 /home/node/clawd/ingest_iatf_docs_to_qdrant.py
"""

import os
import sys
import json
import hashlib
import requests
import traceback
from pathlib import Path
from datetime import datetime

# ── Configuration ────────────────────────────────────────────────────────────────
IATF_ROOT    = "/home/node/clawd/paperless_consume/IATF_documents"
STATE_FILE   = "/home/node/clawd/ingest_iatf_docs_state.json"
LOG_FILE     = "/home/node/clawd/ingest_iatf_docs.log"

INFINITY_URL = "http://infinity:7997/embeddings"
QDRANT_URL   = "http://qdrant:6333"
COLLECTION   = "universal_knowledge"
EMBED_MODEL  = "mixedbread-ai/mxbai-embed-large-v1"
EMBED_DIM    = 1024

CHUNK_SIZE   = 800
ROWS_PER_CHUNK = 8   # XLSX: 何行ずつチャンクにするか
EMBED_BATCH  = 8
QDRANT_BATCH = 32

# ── Logging ──────────────────────────────────────────────────────────────────────
def log(msg, level="INFO"):
    ts   = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    line = f"[{ts}] [{level}] {msg}"
    if sys.stdout.isatty():
        print(line, flush=True)
    else:
        try:
            with open(LOG_FILE, "a", encoding="utf-8") as f:
                f.write(line + "\n")
        except Exception:
            pass

# ── State ─────────────────────────────────────────────────────────────────────────
def load_state():
    if os.path.exists(STATE_FILE):
        try:
            with open(STATE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"processed": {}}

def save_state(state):
    with open(STATE_FILE, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

# ── PPTX parsing ─────────────────────────────────────────────────────────────────
def extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(str(filepath))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [shape.text.strip()
                 for shape in slide.shapes
                 if hasattr(shape, "text") and shape.text.strip()]
        if parts:
            slides.append(f"[スライド{i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)

def make_pptx_chunks(filepath):
    stem = Path(filepath).stem
    try:
        text = extract_pptx(filepath)
    except Exception as e:
        log(f"  PPTX parse error {Path(filepath).name}: {e}", "WARN")
        return []
    if len(text.strip()) < 20:
        return []

    base = {
        "source":   f"iatf_docs/{Path(filepath).name}",
        "category": "IATF16949_解説",
        "filename": stem,
        "filetype": "pptx",
    }
    result = []
    for ci, start in enumerate(range(0, len(text), CHUNK_SIZE)):
        chunk = text[start:start + CHUNK_SIZE]
        if not chunk.strip():
            continue
        pid = int(hashlib.md5(f"iatf_docs::{filepath}::c{ci}".encode()).hexdigest()[:15], 16)
        result.append((chunk, {**base, "chunk": ci, "content": chunk}, pid))
    return result

# ── XLSX parsing ─────────────────────────────────────────────────────────────────
def xlsx_row_to_text(headers, row):
    """ヘッダーと行値をペアにしてテキスト化。空セルはスキップ。"""
    parts = []
    for h, v in zip(headers, row):
        if v is not None and str(v).strip():
            h_str = str(h).strip() if h else ""
            v_str = str(v).strip()
            if h_str:
                parts.append(f"{h_str}: {v_str}")
            else:
                parts.append(v_str)
    return " | ".join(parts)

def make_xlsx_chunks(filepath):
    import openpyxl
    stem  = Path(filepath).stem
    chunks_all = []
    try:
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = [row for row in ws.iter_rows(values_only=True)
                        if any(c is not None and str(c).strip() for c in row)]
            if not all_rows:
                continue

            # 1行目をヘッダーとして扱う
            headers = [str(c).strip() if c else "" for c in all_rows[0]]
            data_rows = all_rows[1:] if len(all_rows) > 1 else []

            if not data_rows:
                # ヘッダーのみの場合はヘッダーテキストだけ1チャンク
                text = f"[ファイル: {stem}][シート: {sheet_name}]\n" + " | ".join(headers)
                chunks_all.append((text, sheet_name, 0))
                continue

            # ROWS_PER_CHUNK 行ずつチャンク
            for ci, start in enumerate(range(0, len(data_rows), ROWS_PER_CHUNK)):
                batch = data_rows[start:start + ROWS_PER_CHUNK]
                lines = [f"[ファイル: {stem}][シート: {sheet_name}]"]
                lines.append("ヘッダー: " + " | ".join(h for h in headers if h))
                for row in batch:
                    row_text = xlsx_row_to_text(headers, row)
                    if row_text:
                        lines.append(row_text)
                text = "\n".join(lines)
                if len(text.strip()) >= 20:
                    chunks_all.append((text, sheet_name, ci))
        wb.close()
    except Exception as e:
        log(f"  XLSX parse error {Path(filepath).name}: {e}", "WARN")
        return []

    base = {
        "source":   f"iatf_docs/{Path(filepath).name}",
        "category": "IATF16949_文書",
        "filename": stem,
        "filetype": "xlsx",
    }
    result = []
    for (text, sheet_name, ci) in chunks_all:
        pid = int(hashlib.md5(f"iatf_docs::{filepath}::{sheet_name}::c{ci}".encode()).hexdigest()[:15], 16)
        result.append((text, {**base, "sheet": sheet_name, "chunk": ci, "content": text[:400]}, pid))
    return result

# ── Infinity batch embed ─────────────────────────────────────────────────────────
def batch_embed(texts):
    try:
        resp = requests.post(INFINITY_URL,
            json={"model": EMBED_MODEL, "input": texts}, timeout=120)
        resp.raise_for_status()
        return [d["embedding"] for d in resp.json()["data"]]
    except Exception as e:
        log(f"  Embed batch error: {e} — 1件ずつリトライ", "WARN")
        results = []
        for t in texts:
            try:
                r = requests.post(INFINITY_URL,
                    json={"model": EMBED_MODEL, "input": t}, timeout=120)
                r.raise_for_status()
                results.append(r.json()["data"][0]["embedding"])
            except Exception as e2:
                log(f"  Embed single error: {e2}", "WARN")
                results.append(None)
        return results

# ── Qdrant batch upsert ──────────────────────────────────────────────────────────
def batch_upsert(points):
    try:
        resp = requests.put(f"{QDRANT_URL}/collections/{COLLECTION}/points",
            json={"points": points}, timeout=30)
        resp.raise_for_status()
        return len(points)
    except Exception as e:
        log(f"  Qdrant error: {e}", "WARN")
        return 0

# ── Core flush ───────────────────────────────────────────────────────────────────
def flush_buffer(chunk_buf):
    if not chunk_buf:
        return
    texts  = [c[0] for c in chunk_buf]
    vecs   = batch_embed(texts)
    points = []
    for (text, payload, pid), vec in zip(chunk_buf, vecs):
        if vec and len(vec) == EMBED_DIM:
            points.append({"id": pid, "vector": vec,
                "payload": {**payload, "ingested_at": datetime.now().isoformat()}})
    for i in range(0, len(points), QDRANT_BATCH):
        batch_upsert(points[i:i + QDRANT_BATCH])
    chunk_buf.clear()

# ── Main ─────────────────────────────────────────────────────────────────────────
def main():
    log("=" * 60)
    log("  ingest_iatf_docs_to_qdrant.py 開始")
    log(f"  ソース: {IATF_ROOT}/")
    log("=" * 60)

    state     = load_state()
    processed = state.get("processed", {})

    root = Path(IATF_ROOT)
    all_files = sorted(list(root.glob("*.pptx")) + list(root.glob("*.xlsx")))
    pending   = [p for p in all_files if str(p) not in processed]

    pptx_cnt = sum(1 for p in pending if p.suffix == ".pptx")
    xlsx_cnt = sum(1 for p in pending if p.suffix == ".xlsx")
    log(f"  未処理: {len(pending)}件 (PPTX={pptx_cnt}  XLSX={xlsx_cnt})")

    if not pending:
        log("  ✅ 新規ファイルなし。")
        return

    ok = skip = err = 0
    chunk_buf = []

    for i, path in enumerate(pending):
        try:
            if path.suffix == ".pptx":
                chunks = make_pptx_chunks(str(path))
                cat = "PPTX"
            elif path.suffix == ".xlsx":
                chunks = make_xlsx_chunks(str(path))
                cat = "XLSX"
            else:
                skip += 1
                continue

            if not chunks:
                processed[str(path)] = {"chunks": 0, "ts": datetime.now().isoformat()}
                skip += 1
                continue

            for item in chunks:
                chunk_buf.append(item)

            if len(chunk_buf) >= EMBED_BATCH:
                flush_buffer(chunk_buf)

            processed[str(path)] = {
                "chunks": len(chunks),
                "type":   cat,
                "ts":     datetime.now().isoformat(),
            }
            ok += 1

        except Exception as e:
            log(f"  [{i+1}] ERROR {path.name}: {e}", "ERROR")
            log(traceback.format_exc(), "ERROR")
            processed[str(path)] = {"error": str(e), "ts": datetime.now().isoformat()}
            err += 1

        if (i + 1) % 20 == 0:
            flush_buffer(chunk_buf)
            state["processed"] = processed
            save_state(state)
            log(f"  [{i+1}/{len(pending)}]  ok={ok}  skip={skip}  err={err}")

    flush_buffer(chunk_buf)
    state["processed"] = processed
    save_state(state)
    log(f"  ✅ 完了: ok={ok}  skip={skip}  err={err}  合計={len(pending)}")

if __name__ == "__main__":
    main()
